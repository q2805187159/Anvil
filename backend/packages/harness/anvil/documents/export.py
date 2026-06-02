from __future__ import annotations

from dataclasses import dataclass
import hashlib
from pathlib import Path
import re
import shutil
from uuid import uuid4

from .contracts import ExportedDocumentResult
from .presentation import build_presentation_contract, build_presentation_plan, preflight_presentation_artifact


SLIDE_WIDTH = 13.333333
SLIDE_HEIGHT = 7.5


@dataclass(frozen=True)
class PresentationVisualProfile:
    title: str
    subtitle: str
    short_title: str
    chips: tuple[str, ...]
    orbit_center: str
    orbit_labels: tuple[str, ...]
    flow_labels: tuple[str, ...]
    metric_labels: tuple[str, ...]
    footer_label: str
    seed: int


def export_document(
    *,
    output_path: Path,
    content: str,
    format: str = "docx",
    mode: str = "editable",
    scratch_root: Path | None = None,
    cleanup_intermediates: bool = True,
) -> ExportedDocumentResult:
    normalized_format = format.strip().lower()
    if normalized_format not in {"docx", "pptx"}:
        raise ValueError(f"unsupported export format: {format}")

    normalized_mode = mode.strip().lower()
    warnings: list[str] = []
    effective_mode = normalized_mode
    if normalized_mode == "preserve_layout":
        warnings.append("preserve_layout provider is unavailable; falling back to editable mode.")
        effective_mode = "editable"

    scratch_root = scratch_root or (output_path.parent.parent / "workspace" / ".anvil-scratch")
    scratch_dir = scratch_root / f"export-{uuid4().hex[:12]}"
    scratch_dir.mkdir(parents=True, exist_ok=True)
    scratch_md_path = scratch_dir / f"{output_path.stem}.md"
    scratch_md_path.write_text(content, encoding="utf-8")

    output_path.parent.mkdir(parents=True, exist_ok=True)
    if normalized_format == "docx":
        _write_docx_from_markdown(scratch_md_path, output_path)
        provider = "python-docx"
        preflight = None
        metadata: dict[str, object] = {}
    else:
        metadata = _write_pptx_from_markdown(scratch_md_path, output_path)
        provider = "python-pptx"
        contract = build_presentation_contract(content=content, title=output_path.stem, max_slides=40)
        preflight = preflight_presentation_artifact(output_path, contract=contract).to_payload()

    cleaned_scratch_paths: list[Path] = []
    if cleanup_intermediates:
        shutil.rmtree(scratch_dir, ignore_errors=True)
        cleaned_scratch_paths.append(scratch_dir)

    return ExportedDocumentResult(
        output_path=output_path,
        mode=effective_mode,
        format=normalized_format,
        provider=provider,
        warnings=tuple(warnings),
        scratch_paths=(scratch_md_path,),
        cleaned_scratch_paths=tuple(cleaned_scratch_paths),
        preflight=preflight,
        metadata=metadata,
    )


def _write_docx_from_markdown(markdown_path: Path, output_path: Path) -> None:
    from docx import Document

    document = Document()
    lines = markdown_path.read_text(encoding="utf-8").splitlines()
    for line in lines:
        stripped = line.strip()
        if not stripped:
            continue
        if stripped.startswith("#"):
            level = len(stripped) - len(stripped.lstrip("#"))
            text = stripped[level:].strip()
            if level == 1:
                paragraph = document.add_paragraph()
                paragraph.style = "Title"
                paragraph.add_run(text)
            else:
                document.add_heading(text, level=min(level - 1, 4))
            continue
        if stripped.startswith(">"):
            paragraph = document.add_paragraph(style="Intense Quote")
            paragraph.add_run(stripped.lstrip(">").strip())
            continue
        if stripped.startswith(("- ", "* ", "• ")):
            document.add_paragraph(stripped[2:].strip(), style="List Bullet")
            continue
        document.add_paragraph(stripped)
    document.save(str(output_path))


def _build_visual_profile(content: str, title: str) -> PresentationVisualProfile:
    title = title.strip() or "Presentation"
    terms = _extract_topic_terms(content, title)
    chips = tuple(_pad_terms(terms[:4], title, ("Scope", "Flow", "Evidence", "Next")))
    orbit_labels = tuple(_pad_terms(terms[1:5], title, ("Input", "System", "Signals", "Outcome")))
    flow_labels = tuple(_pad_terms(terms[:5], title, ("Input", "Context", "Policy", "Tools", "Result")))
    metric_labels = tuple(_pad_terms(terms[2:5], title, ("Readiness", "Clarity", "Quality")))
    subtitle = _extract_subtitle(content, title)
    short_title = _short_label(title)
    seed = int(hashlib.sha256(f"{title}\n{content[:4000]}".encode("utf-8", errors="replace")).hexdigest()[:8], 16)
    return PresentationVisualProfile(
        title=title,
        subtitle=subtitle,
        short_title=short_title,
        chips=chips,
        orbit_center=short_title,
        orbit_labels=orbit_labels,
        flow_labels=flow_labels,
        metric_labels=metric_labels,
        footer_label=short_title,
        seed=seed,
    )


def _extract_topic_terms(content: str, title: str) -> list[str]:
    candidates: list[str] = []
    for line in content.splitlines():
        stripped = line.strip().strip("#-*•>0123456789. ")
        if not stripped:
            continue
        if line.lstrip().startswith("#"):
            candidates.extend(_split_term_candidates(stripped))
            continue
        if line.lstrip().startswith(("- ", "* ", "• ")):
            candidates.extend(_split_term_candidates(stripped.split(":", 1)[0]))
    candidates.extend(_split_term_candidates(title))
    seen: set[str] = set()
    terms: list[str] = []
    for candidate in candidates:
        normalized = candidate.strip(" -_/:：|")
        if len(normalized) < 2:
            continue
        key = normalized.casefold()
        if key in seen or key in _COMMON_TOPIC_WORDS:
            continue
        seen.add(key)
        terms.append(_short_label(normalized))
        if len(terms) >= 12:
            break
    return terms


_COMMON_TOPIC_WORDS = {
    "about",
    "agent",
    "agents",
    "deck",
    "introduction",
    "overview",
    "presentation",
    "project",
    "slide",
    "slides",
    "介绍",
    "项目",
}


def _split_term_candidates(value: str) -> list[str]:
    if re.search(r"[\u4e00-\u9fff]", value):
        pieces = re.split(r"[，,、/|：:（）()【】\[\]\s]+", value)
        return [piece for piece in pieces if piece]
    words = re.findall(r"[A-Za-z][A-Za-z0-9+.-]{2,}", value)
    phrases = re.split(r"[,/|:;()\[\]-]+", value)
    return [*phrases[:2], *words]


def _pad_terms(terms: list[str], title: str, fallback: tuple[str, ...]) -> list[str]:
    result = [term for term in terms if term]
    for term in (_short_label(title), *fallback):
        if len(result) >= len(fallback):
            break
        if term and term.casefold() not in {item.casefold() for item in result}:
            result.append(term)
    return result


def _extract_subtitle(content: str, title: str) -> str:
    for line in content.splitlines():
        stripped = line.strip()
        if not stripped or stripped.startswith("#") or stripped.startswith(("- ", "* ", "• ")):
            continue
        return _clean_bullet(stripped)
    terms = _extract_topic_terms(content, title)
    if len(terms) >= 3:
        return f"Focused view of {terms[0]}, {terms[1]}, and {terms[2]}."
    return f"A polished editable deck about {_short_label(title)}."


def _short_label(value: str, *, max_chars: int = 18) -> str:
    cleaned = value.replace("**", "").replace("__", "").strip()
    if not cleaned:
        return "Deck"
    if len(cleaned) <= max_chars:
        return cleaned
    if re.search(r"[\u4e00-\u9fff]", cleaned):
        return cleaned[:max_chars]
    words = cleaned.split()
    label = ""
    for word in words:
        next_label = f"{label} {word}".strip()
        if len(next_label) > max_chars:
            break
        label = next_label
    return label or cleaned[:max_chars].rstrip()


def _palette_for_seed(seed: int):
    from pptx.dml.color import RGBColor

    palettes = [
        {
            "navy": RGBColor(18, 27, 44),
            "blue": RGBColor(42, 122, 194),
            "cyan": RGBColor(34, 197, 219),
            "green": RGBColor(76, 175, 118),
            "amber": RGBColor(245, 174, 70),
        },
        {
            "navy": RGBColor(31, 35, 52),
            "blue": RGBColor(63, 118, 176),
            "cyan": RGBColor(42, 180, 171),
            "green": RGBColor(103, 162, 90),
            "amber": RGBColor(222, 153, 74),
        },
        {
            "navy": RGBColor(26, 40, 48),
            "blue": RGBColor(56, 111, 164),
            "cyan": RGBColor(69, 181, 202),
            "green": RGBColor(86, 154, 119),
            "amber": RGBColor(218, 168, 72),
        },
    ]
    return palettes[seed % len(palettes)]


def _write_pptx_from_markdown(markdown_path: Path, output_path: Path) -> dict[str, object]:
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt
    except ImportError as exc:
        raise ValueError("pptx export requires python-pptx") from exc

    content = markdown_path.read_text(encoding="utf-8")
    plan = build_presentation_plan(content=content, title=output_path.stem, max_slides=40)
    profile = _build_visual_profile(content, plan.title)
    presentation = Presentation()
    presentation.slide_width = Inches(SLIDE_WIDTH)
    presentation.slide_height = Inches(SLIDE_HEIGHT)
    blank_layout = presentation.slide_layouts[6]

    palette = _palette_for_seed(profile.seed)
    theme = {
        "navy": palette["navy"],
        "ink": RGBColor(27, 36, 50),
        "muted": RGBColor(88, 100, 117),
        "blue": palette["blue"],
        "cyan": palette["cyan"],
        "green": palette["green"],
        "amber": palette["amber"],
        "red": RGBColor(228, 92, 92),
        "paper": RGBColor(246, 249, 252),
        "white": RGBColor(255, 255, 255),
        "line": RGBColor(217, 226, 236),
    }
    embedded_visuals = _build_embedded_visual_assets(markdown_path.parent, profile=profile, theme=theme)
    visual_assets = 0
    picture_assets = 0

    title_slide = presentation.slides.add_slide(blank_layout)
    _add_background(title_slide, theme, dark=True)
    if "network" in embedded_visuals:
        picture_assets += _add_picture(title_slide, embedded_visuals["network"], 8.12, 0.72, 4.45, 4.86)
    _add_title(title_slide, plan.title, x=0.75, y=0.82, width=7.3, size=34, color=theme["white"])
    _add_text(
        title_slide,
        profile.subtitle,
        x=0.78,
        y=2.0,
        width=6.7,
        height=0.9,
        size=18,
        color=RGBColor(200, 215, 232),
    )
    for pill_index, chip in enumerate(profile.chips[:4]):
        _add_pill(title_slide, chip.upper(), 0.8 + pill_index * 1.32, 3.18, [theme["cyan"], theme["green"], theme["amber"], theme["blue"]][pill_index % 4], theme)
    visual_assets += _add_system_orbit(title_slide, theme, profile)
    _add_footer(title_slide, 1, len(plan.slides) + 1, theme, profile=profile, dark=True)

    for item in plan.slides:
        index = int(item.get("index") or 1)
        title = str(item.get("title") or f"Slide {index}")
        bullets = [str(bullet) for bullet in list(item.get("bullets") or []) if str(bullet).strip()]
        if not bullets:
            bullets = ["Define the claim, proof object, and next action for this slide."]
        slide = presentation.slides.add_slide(blank_layout)
        _add_background(slide, theme, dark=False)
        _add_section_marker(slide, index + 1, theme)
        _add_title(slide, title, x=0.65, y=0.48, width=7.8, size=25, color=theme["ink"])
        layout = _layout_variant(profile, index)
        if layout == 0:
            visual_assets += _add_card_grid(slide, bullets, theme)
            visual_assets += _add_metric_stack(slide, bullets, theme, profile)
            if "matrix" in embedded_visuals:
                picture_assets += _add_picture(slide, embedded_visuals["matrix"], 9.08, 4.66, 2.65, 1.18)
        elif layout == 1:
            visual_assets += _add_architecture_flow(slide, bullets, theme, profile)
            visual_assets += _add_sidebar_notes(slide, bullets, theme)
            if "network_light" in embedded_visuals:
                picture_assets += _add_picture(slide, embedded_visuals["network_light"], 8.96, 4.64, 2.8, 1.25)
        elif layout == 2:
            visual_assets += _add_timeline(slide, bullets, theme)
            visual_assets += _add_quote_panel(slide, bullets, theme)
            if "signals" in embedded_visuals:
                picture_assets += _add_picture(slide, embedded_visuals["signals"], 9.1, 4.53, 2.55, 0.92)
        elif layout == 3:
            visual_assets += _add_split_proof(slide, bullets, theme)
            visual_assets += _add_signal_bars(slide, bullets, theme)
            if "matrix" in embedded_visuals:
                picture_assets += _add_picture(slide, embedded_visuals["matrix"], 4.58, 4.56, 3.08, 1.02)
        elif layout == 4:
            visual_assets += _add_radar_map(slide, bullets, theme, profile)
            visual_assets += _add_sidebar_notes(slide, bullets[1:] or bullets, theme)
            if "signals" in embedded_visuals:
                picture_assets += _add_picture(slide, embedded_visuals["signals"], 0.9, 5.0, 3.8, 0.9)
        else:
            visual_assets += _add_comparison_lanes(slide, bullets, theme, profile)
            if "network_light" in embedded_visuals:
                picture_assets += _add_picture(slide, embedded_visuals["network_light"], 9.3, 1.48, 2.55, 1.3)
        _add_footer(slide, index + 1, len(plan.slides) + 1, theme, profile=profile, dark=False)

    presentation.save(str(output_path))
    return {
        "visual_system": "editorial-topic-aware",
        "slide_count": len(presentation.slides),
        "visual_asset_count": visual_assets,
        "embedded_picture_count": picture_assets,
        "topic_profile": {
            "title": profile.title,
            "chips": list(profile.chips),
            "footer_label": profile.footer_label,
            "seed": profile.seed,
            "layout_family": profile.seed % 6,
        },
        "scratch_policy": "intermediates stay in thread scratch; final output path receives only the deck",
    }


def _layout_variant(profile: PresentationVisualProfile, index: int) -> int:
    return (index + (profile.seed % 6)) % 6


def _add_background(slide, theme, *, dark: bool) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(SLIDE_WIDTH), Inches(SLIDE_HEIGHT))
    bg.fill.solid()
    bg.fill.fore_color.rgb = theme["navy"] if dark else theme["paper"]
    bg.line.fill.background()
    if dark:
        accent = slide.shapes.add_shape(MSO_SHAPE.ARC, Inches(8.4), Inches(-0.6), Inches(5.8), Inches(5.8))
        accent.line.color.rgb = theme["cyan"]
        accent.line.width = Inches(0.03)
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.78), Inches(SLIDE_WIDTH), Inches(0.72))
        band.fill.solid()
        band.fill.fore_color.rgb = RGBColor(10, 18, 32)
        band.line.fill.background()
        return
    rail = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), Inches(SLIDE_HEIGHT))
    rail.fill.solid()
    rail.fill.fore_color.rgb = theme["blue"]
    rail.line.fill.background()
    wash = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.45), Inches(0), Inches(3.88), Inches(SLIDE_HEIGHT))
    wash.fill.solid()
    wash.fill.fore_color.rgb = RGBColor(235, 243, 250)
    wash.line.fill.background()


def _add_title(slide, text: str, *, x: float, y: float, width: float, size: int, color) -> None:
    from pptx.util import Inches, Pt

    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(0.9))
    frame = box.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text[:110]
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.name = "Microsoft YaHei"
    run.font.color.rgb = color


def _add_text(slide, text: str, *, x: float, y: float, width: float, height: float, size: int, color, bold: bool = False) -> None:
    from pptx.util import Inches, Pt

    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.line_spacing = 1.12
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_pill(slide, text: str, x: float, y: float, color, theme) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(1.08), Inches(0.36))
    pill.fill.solid()
    pill.fill.fore_color.rgb = color
    pill.line.fill.background()
    frame = pill.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = theme["navy"]


def _add_system_orbit(slide, theme, profile: PresentationVisualProfile) -> int:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    labels = [
        (profile.orbit_labels[0], 9.25, 1.25, theme["cyan"]),
        (profile.orbit_labels[1], 10.45, 2.6, theme["blue"]),
        (profile.orbit_labels[2], 9.05, 4.0, theme["green"]),
        (profile.orbit_labels[3], 11.05, 4.4, theme["amber"]),
    ]
    center = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.7), Inches(2.55), Inches(1.65), Inches(1.65))
    center.fill.solid()
    center.fill.fore_color.rgb = theme["white"]
    center.line.color.rgb = theme["cyan"]
    center.line.width = Inches(0.02)
    center.text_frame.clear()
    center.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = center.text_frame.paragraphs[0].add_run()
    run.text = profile.orbit_center
    run.font.name = "Arial"
    run.font.bold = True
    run.font.size = Pt(14)
    run.font.color.rgb = theme["navy"]
    count = 1
    for label, x, y, color in labels:
        node = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(1.55), Inches(0.48))
        node.fill.solid()
        node.fill.fore_color.rgb = color
        node.line.fill.background()
        node.text_frame.clear()
        paragraph = node.text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        run = paragraph.add_run()
        run.text = label
        run.font.name = "Arial"
        run.font.bold = True
        run.font.size = Pt(10)
        run.font.color.rgb = theme["navy"]
        count += 1
    return count


def _add_section_marker(slide, index: int, theme) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    marker = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.55), Inches(6.58), Inches(0.48), Inches(0.48))
    marker.fill.solid()
    marker.fill.fore_color.rgb = theme["ink"]
    marker.line.fill.background()
    paragraph = marker.text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = f"{index:02d}"
    run.font.name = "Arial"
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = theme["white"]


def _add_card_grid(slide, bullets: list[str], theme) -> int:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    cards = bullets[:4]
    count = 0
    for idx, bullet in enumerate(cards):
        row = idx // 2
        col = idx % 2
        x = 0.72 + col * 3.45
        y = 1.78 + row * 1.58
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.02), Inches(1.15))
        card.fill.solid()
        card.fill.fore_color.rgb = theme["white"]
        card.line.color.rgb = theme["line"]
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.12), Inches(1.15))
        bar.fill.solid()
        bar.fill.fore_color.rgb = [theme["cyan"], theme["green"], theme["amber"], theme["blue"]][idx]
        bar.line.fill.background()
        _add_text(slide, _clean_bullet(bullet), x=x + 0.24, y=y + 0.16, width=2.58, height=0.78, size=12, color=theme["ink"], bold=True)
        count += 2
    return count


def _add_metric_stack(slide, bullets: list[str], theme, profile: PresentationVisualProfile) -> int:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    labels = list(profile.metric_labels[:3])
    count = 0
    for idx, label in enumerate(labels):
        y = 1.85 + idx * 1.18
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.15), Inches(y), Inches(2.8), Inches(0.78))
        shape.fill.solid()
        shape.fill.fore_color.rgb = [theme["blue"], theme["green"], theme["amber"]][idx]
        shape.line.fill.background()
        paragraph = shape.text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        run = paragraph.add_run()
        run.text = label
        run.font.name = "Arial"
        run.font.size = Pt(13)
        run.font.bold = True
        run.font.color.rgb = theme["white"]
        count += 1
    return count


def _add_architecture_flow(slide, bullets: list[str], theme, profile: PresentationVisualProfile) -> int:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    labels = list(profile.flow_labels[:5])
    count = 0
    for idx, label in enumerate(labels):
        x = 0.78 + idx * 1.58
        node = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.55), Inches(1.18), Inches(0.72))
        node.fill.solid()
        node.fill.fore_color.rgb = [theme["blue"], theme["cyan"], theme["amber"], theme["green"], theme["ink"]][idx]
        node.line.fill.background()
        paragraph = node.text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        run = paragraph.add_run()
        run.text = label
        run.font.name = "Arial"
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = theme["white"]
        count += 1
        if idx < len(labels) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 1.16), Inches(2.73), Inches(0.5), Inches(0.24))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = theme["line"]
            arrow.line.fill.background()
            count += 1
    return count


def _add_sidebar_notes(slide, bullets: list[str], theme) -> int:
    notes = bullets[:3]
    for idx, note in enumerate(notes):
        _add_text(slide, f"{idx + 1}. {_clean_bullet(note)}", x=9.05, y=1.78 + idx * 1.15, width=2.9, height=0.78, size=12, color=theme["ink"])
    return len(notes)


def _add_timeline(slide, bullets: list[str], theme) -> int:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    items = bullets[:4]
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.95), Inches(3.05), Inches(6.5), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = theme["line"]
    line.line.fill.background()
    count = 1
    for idx, item in enumerate(items):
        x = 0.9 + idx * 1.7
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(2.82), Inches(0.5), Inches(0.5))
        dot.fill.solid()
        dot.fill.fore_color.rgb = [theme["cyan"], theme["blue"], theme["green"], theme["amber"]][idx % 4]
        dot.line.fill.background()
        _add_text(slide, _clean_bullet(item), x=x - 0.08, y=3.48, width=1.5, height=0.95, size=10, color=theme["ink"], bold=True)
        count += 2
    return count


def _add_quote_panel(slide, bullets: list[str], theme) -> int:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.85), Inches(1.55), Inches(3.15), Inches(3.7))
    panel.fill.solid()
    panel.fill.fore_color.rgb = theme["navy"]
    panel.line.fill.background()
    text = _clean_bullet(bullets[0] if bullets else "Harness keeps execution observable and governed.")
    _add_text(slide, text, x=9.18, y=2.08, width=2.46, height=2.15, size=17, color=theme["white"], bold=True)
    return 1


def _add_split_proof(slide, bullets: list[str], theme) -> int:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    left = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.78), Inches(1.65), Inches(3.25), Inches(3.88))
    left.fill.solid()
    left.fill.fore_color.rgb = theme["ink"]
    left.line.fill.background()
    _add_text(slide, "CLAIM", x=1.04, y=1.92, width=1.8, height=0.35, size=10, color=theme["cyan"], bold=True)
    _add_text(slide, _clean_bullet(bullets[0] if bullets else "One clear operating claim."), x=1.04, y=2.38, width=2.54, height=1.45, size=18, color=theme["white"], bold=True)
    _add_text(slide, _clean_bullet(bullets[1] if len(bullets) > 1 else "Proof stays attached to tools and state."), x=4.55, y=1.88, width=3.2, height=0.88, size=13, color=theme["ink"], bold=True)
    _add_text(slide, _clean_bullet(bullets[2] if len(bullets) > 2 else "The final answer is only one surface of the run."), x=4.55, y=3.08, width=3.2, height=1.0, size=13, color=theme["muted"])
    return 3


def _add_signal_bars(slide, bullets: list[str], theme) -> int:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    values = [0.82, 0.66, 0.74, 0.58]
    count = 0
    for idx, value in enumerate(values):
        y = 1.85 + idx * 0.78
        track = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.9), Inches(y), Inches(2.75), Inches(0.26))
        track.fill.solid()
        track.fill.fore_color.rgb = theme["line"]
        track.line.fill.background()
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.9), Inches(y), Inches(2.75 * value), Inches(0.26))
        bar.fill.solid()
        bar.fill.fore_color.rgb = [theme["blue"], theme["cyan"], theme["green"], theme["amber"]][idx]
        bar.line.fill.background()
        count += 2
    return count


def _add_radar_map(slide, bullets: list[str], theme, profile: PresentationVisualProfile) -> int:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    labels = list(_pad_terms([_clean_bullet(item).split(" ", 1)[0] for item in bullets[:5]], profile.title, ("Input", "Signal", "Policy", "Output", "Review")))
    center_x, center_y = 4.28, 3.35
    radii = [2.45, 1.76, 1.08]
    count = 0
    for idx, radius in enumerate(radii):
        ring = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(center_x - radius), Inches(center_y - radius), Inches(radius * 2), Inches(radius * 2))
        ring.fill.background()
        ring.line.color.rgb = [theme["line"], theme["cyan"], theme["blue"]][idx]
        ring.line.width = Inches(0.014)
        count += 1
    positions = [(2.1, 1.45), (5.35, 1.35), (6.48, 3.4), (4.52, 5.32), (1.66, 4.42)]
    colors = [theme["blue"], theme["cyan"], theme["green"], theme["amber"], theme["ink"]]
    for idx, label in enumerate(labels[:5]):
        x, y = positions[idx]
        node = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(1.58), Inches(0.54))
        node.fill.solid()
        node.fill.fore_color.rgb = colors[idx]
        node.line.fill.background()
        paragraph = node.text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        run = paragraph.add_run()
        run.text = label
        run.font.name = "Arial"
        run.font.size = Pt(10)
        run.font.bold = True
        run.font.color.rgb = theme["white"]
        count += 1
    _add_text(slide, profile.short_title, x=center_x - 0.68, y=center_y - 0.18, width=1.38, height=0.38, size=12, color=theme["ink"], bold=True)
    return count + 1


def _add_comparison_lanes(slide, bullets: list[str], theme, profile: PresentationVisualProfile) -> int:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    left_title = profile.flow_labels[0] if profile.flow_labels else "Current"
    right_title = profile.flow_labels[-1] if profile.flow_labels else "Target"
    lanes = [
        (0.84, theme["navy"], left_title, bullets[:3]),
        (6.0, theme["blue"], right_title, bullets[3:6] or bullets[:3]),
    ]
    count = 0
    for x, color, title, lane_bullets in lanes:
        panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.55), Inches(4.45), Inches(4.42))
        panel.fill.solid()
        panel.fill.fore_color.rgb = theme["white"]
        panel.line.color.rgb = theme["line"]
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.55), Inches(4.45), Inches(0.52))
        header.fill.solid()
        header.fill.fore_color.rgb = color
        header.line.fill.background()
        _add_text(slide, title, x=x + 0.22, y=1.66, width=3.8, height=0.28, size=11, color=theme["white"], bold=True)
        for idx, bullet in enumerate(lane_bullets[:3]):
            _add_text(slide, f"{idx + 1}. {_clean_bullet(bullet)}", x=x + 0.28, y=2.38 + idx * 0.92, width=3.76, height=0.58, size=12, color=theme["ink"])
        count += 2 + len(lane_bullets[:3])
    bridge = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(5.16), Inches(3.32), Inches(0.58), Inches(0.36))
    bridge.fill.solid()
    bridge.fill.fore_color.rgb = theme["amber"]
    bridge.line.fill.background()
    return count + 1


def _add_footer(slide, index: int, total: int, theme, *, profile: PresentationVisualProfile, dark: bool) -> None:
    from pptx.dml.color import RGBColor

    color = RGBColor(175, 190, 210) if dark else theme["muted"]
    _add_text(slide, f"{profile.footer_label} / {index:02d}-{total:02d}", x=10.25, y=6.86, width=2.05, height=0.28, size=8, color=color)


def _clean_bullet(value: str) -> str:
    cleaned = value.replace("**", "").replace("__", "").strip()
    if len(cleaned) > 92:
        return cleaned[:89].rstrip() + "..."
    return cleaned


def _add_picture(slide, image_path: Path, x: float, y: float, width: float, height: float) -> int:
    from pptx.util import Inches

    if not image_path.exists():
        return 0
    slide.shapes.add_picture(str(image_path), Inches(x), Inches(y), width=Inches(width), height=Inches(height))
    return 1


def _build_embedded_visual_assets(scratch_dir: Path, *, profile: PresentationVisualProfile, theme) -> dict[str, Path]:
    try:
        from PIL import Image, ImageDraw
    except ImportError:
        return {}
    asset_dir = scratch_dir / "ppt-assets"
    asset_dir.mkdir(parents=True, exist_ok=True)

    seed = profile.seed

    def color_tuple(value, alpha: int) -> tuple[int, int, int, int]:
        return (int(value[0]), int(value[1]), int(value[2]), alpha)

    def jitter(value: int, index: int, span: int) -> int:
        return value + ((seed >> (index % 16)) % (span * 2 + 1)) - span

    def save_network(path: Path, *, dark: bool) -> None:
        bg = (18, 27, 44, 0 if dark else 255)
        image = Image.new("RGBA", (900, 620), bg)
        draw = ImageDraw.Draw(image, "RGBA")
        base_nodes = [(130, 130), (420, 95), (710, 160), (280, 360), (610, 410), (420, 260)]
        nodes = [(jitter(x, idx * 2, 28), jitter(y, idx * 2 + 1, 22)) for idx, (x, y) in enumerate(base_nodes)]
        line = color_tuple(theme["cyan"], 120) if dark else color_tuple(theme["blue"], 105)
        fill = (255, 255, 255, 215) if dark else color_tuple(theme["blue"], 225)
        for start, end in [(0, 1), (1, 2), (0, 3), (3, 5), (5, 2), (5, 4), (3, 4)]:
            draw.line((nodes[start], nodes[end]), fill=line, width=9)
        for idx, (x, y) in enumerate(nodes):
            radius = 40 if idx != 5 else 58
            color = fill if idx != 5 else (color_tuple(theme["cyan"], 235) if dark else color_tuple(theme["navy"], 235))
            draw.ellipse((x - radius, y - radius, x + radius, y + radius), fill=color, outline=(255, 255, 255, 180), width=4)
        image.save(path)

    def save_matrix(path: Path) -> None:
        image = Image.new("RGBA", (760, 280), (255, 255, 255, 0))
        draw = ImageDraw.Draw(image, "RGBA")
        colors = [color_tuple(theme["blue"], 230), color_tuple(theme["cyan"], 230), color_tuple(theme["green"], 230), color_tuple(theme["amber"], 230)]
        for row in range(4):
            for col in range(8):
                x = 28 + col * 86 + ((seed >> ((row + col) % 12)) % 9)
                y = 26 + row * 58 + ((seed >> ((row * 3 + col) % 12)) % 7)
                alpha = 95 + ((row + col + seed) % 4) * 38
                color = colors[(row + col + seed) % len(colors)]
                draw.rounded_rectangle((x, y, x + 58, y + 34), radius=8, fill=(color[0], color[1], color[2], alpha))
        image.save(path)

    def save_signals(path: Path) -> None:
        image = Image.new("RGBA", (720, 240), (255, 255, 255, 0))
        draw = ImageDraw.Draw(image, "RGBA")
        base_points = [(32, 176), (120, 122), (210, 150), (310, 86), (410, 112), (520, 52), (666, 88)]
        points = [(x, jitter(y, idx + 5, 24)) for idx, (x, y) in enumerate(base_points)]
        draw.line(points, fill=color_tuple(theme["blue"], 230), width=10, joint="curve")
        for x, y in points:
            draw.ellipse((x - 14, y - 14, x + 14, y + 14), fill=color_tuple(theme["cyan"], 245), outline=color_tuple(theme["navy"], 180), width=3)
        for idx, height in enumerate([92, 132, 72, 156, 114, 178]):
            x = 70 + idx * 92
            adjusted = max(42, min(188, height + ((seed >> idx) % 39) - 19))
            draw.rounded_rectangle((x, 216 - adjusted, x + 32, 216), radius=7, fill=color_tuple(theme["green"], 120))
        image.save(path)

    suffix = f"{profile.seed:08x}"
    assets = {
        "network": asset_dir / f"topic-network-{suffix}.png",
        "network_light": asset_dir / f"topic-network-light-{suffix}.png",
        "matrix": asset_dir / f"topic-matrix-{suffix}.png",
        "signals": asset_dir / f"topic-signals-{suffix}.png",
    }
    save_network(assets["network"], dark=True)
    save_network(assets["network_light"], dark=False)
    save_matrix(assets["matrix"])
    save_signals(assets["signals"])
    return assets

