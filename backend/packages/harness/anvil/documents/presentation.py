from __future__ import annotations

import base64
from dataclasses import dataclass
import html
import json
from pathlib import Path
import re
import shutil
from xml.etree import ElementTree as ET


MAX_SLIDES = 40
MAX_BULLETS_PER_SLIDE = 8
MAX_TEMPLATE_PREVIEW_PAGES = 16
MAX_TEMPLATE_PREVIEW_SVG_BYTES = 1_500_000
MAX_VISUAL_REVIEW_PAGES = 40
VISUAL_SIGNATURE_GRID_COLUMNS = 8
VISUAL_SIGNATURE_GRID_ROWS = 5


@dataclass(frozen=True)
class PresentationPlan:
    title: str
    slides: tuple[dict[str, object], ...]
    source: str
    warnings: tuple[str, ...] = ()

    def to_payload(self) -> dict[str, object]:
        return {
            "title": self.title,
            "slides": list(self.slides),
            "slide_count": len(self.slides),
            "source": self.source,
            "warnings": list(self.warnings),
        }


@dataclass(frozen=True)
class PresentationContract:
    title: str
    canvas: dict[str, object]
    theme: dict[str, object]
    typography: dict[str, object]
    slides: tuple[dict[str, object], ...]
    assets: tuple[dict[str, object], ...] = ()
    source: str = "presentation_plan"
    warnings: tuple[str, ...] = ()

    def to_payload(self) -> dict[str, object]:
        return {
            "title": self.title,
            "canvas": dict(self.canvas),
            "theme": dict(self.theme),
            "typography": dict(self.typography),
            "slides": list(self.slides),
            "assets": list(self.assets),
            "source": self.source,
            "warnings": list(self.warnings),
        }


@dataclass(frozen=True)
class PresentationPreflightIssue:
    severity: str
    code: str
    message: str
    path: str | None = None

    def to_payload(self) -> dict[str, object]:
        payload: dict[str, object] = {
            "severity": self.severity,
            "code": self.code,
            "message": self.message,
        }
        if self.path is not None:
            payload["path"] = self.path
        return payload


@dataclass(frozen=True)
class PresentationPreflightResult:
    status: str
    format: str
    checked_paths: tuple[str, ...]
    issues: tuple[PresentationPreflightIssue, ...] = ()
    summary: dict[str, object] | None = None

    def to_payload(self) -> dict[str, object]:
        errors = sum(1 for issue in self.issues if issue.severity == "error")
        warnings = sum(1 for issue in self.issues if issue.severity == "warning")
        return {
            "status": self.status,
            "format": self.format,
            "checked_paths": list(self.checked_paths),
            "issues": [issue.to_payload() for issue in self.issues],
            "error_count": errors,
            "warning_count": warnings,
            "summary": dict(self.summary or {}),
        }


@dataclass(frozen=True)
class PresentationTemplateInfo:
    template_id: str
    template_path: Path
    valid: bool
    design_spec_path: Path | None
    pages: tuple[dict[str, object], ...]
    assets: tuple[dict[str, object], ...]
    metadata: dict[str, object]
    warnings: tuple[str, ...] = ()

    def to_payload(self) -> dict[str, object]:
        return {
            "template_id": self.template_id,
            "template_path": str(self.template_path),
            "valid": self.valid,
            "design_spec_path": str(self.design_spec_path) if self.design_spec_path else None,
            "pages": list(self.pages),
            "assets": list(self.assets),
            "metadata": dict(self.metadata),
            "warnings": list(self.warnings),
        }


@dataclass(frozen=True)
class PresentationTemplateImportResult:
    status: str
    template_id: str
    project_path: Path
    template_dir: Path
    copied_paths: tuple[Path, ...]
    skipped_paths: tuple[Path, ...] = ()
    warnings: tuple[str, ...] = ()

    def to_payload(self) -> dict[str, object]:
        return {
            "status": self.status,
            "template_id": self.template_id,
            "project_path": str(self.project_path),
            "template_dir": str(self.template_dir),
            "copied_paths": [str(path) for path in self.copied_paths],
            "skipped_paths": [str(path) for path in self.skipped_paths],
            "warnings": list(self.warnings),
        }


@dataclass(frozen=True)
class PresentationTemplatePreviewResult:
    status: str
    template_id: str
    preview_path: Path
    manifest_path: Path
    pages: tuple[dict[str, object], ...]
    warnings: tuple[str, ...] = ()

    def to_payload(self) -> dict[str, object]:
        return {
            "status": self.status,
            "template_id": self.template_id,
            "preview_path": str(self.preview_path),
            "manifest_path": str(self.manifest_path),
            "page_count": len(self.pages),
            "pages": list(self.pages),
            "warnings": list(self.warnings),
        }


@dataclass(frozen=True)
class PresentationVisualIssue:
    severity: str
    code: str
    message: str

    def to_payload(self) -> dict[str, object]:
        return {
            "severity": self.severity,
            "code": self.code,
            "message": self.message,
        }


@dataclass(frozen=True)
class PresentationVisualPageReview:
    index: int
    name: str
    path: Path
    status: str
    view_box: str | None
    metrics: dict[str, object]
    contract: dict[str, object]
    issues: tuple[PresentationVisualIssue, ...] = ()
    recommendations: tuple[str, ...] = ()
    preflight: dict[str, object] | None = None

    def to_payload(self) -> dict[str, object]:
        payload: dict[str, object] = {
            "index": self.index,
            "name": self.name,
            "path": str(self.path),
            "status": self.status,
            "viewBox": self.view_box,
            "metrics": dict(self.metrics),
            "contract": dict(self.contract),
            "issues": [issue.to_payload() for issue in self.issues],
            "recommendations": list(self.recommendations),
        }
        if self.preflight is not None:
            payload["preflight"] = dict(self.preflight)
        return payload


@dataclass(frozen=True)
class PresentationVisualReviewResult:
    status: str
    target_path: Path
    checked_paths: tuple[Path, ...]
    pages: tuple[PresentationVisualPageReview, ...]
    summary: dict[str, object]
    recommendations: tuple[str, ...] = ()
    truncated: bool = False

    def to_payload(self) -> dict[str, object]:
        return {
            "status": self.status,
            "target_path": str(self.target_path),
            "checked_paths": [str(path) for path in self.checked_paths],
            "checked_count": len(self.checked_paths),
            "truncated": self.truncated,
            "summary": dict(self.summary),
            "pages": [page.to_payload() for page in self.pages],
            "recommendations": list(self.recommendations),
        }


@dataclass(frozen=True)
class PresentationReviewReportResult:
    status: str
    report_path: Path
    manifest_path: Path
    target_path: Path
    checked_paths: tuple[Path, ...]
    summary: dict[str, object]
    recommendations: tuple[str, ...] = ()
    warnings: tuple[str, ...] = ()

    def to_payload(self) -> dict[str, object]:
        return {
            "status": self.status,
            "report_path": str(self.report_path),
            "manifest_path": str(self.manifest_path),
            "target_path": str(self.target_path),
            "checked_paths": [str(path) for path in self.checked_paths],
            "checked_count": len(self.checked_paths),
            "summary": dict(self.summary),
            "recommendations": list(self.recommendations),
            "warnings": list(self.warnings),
        }


@dataclass(frozen=True)
class PresentationVisualDiffPage:
    index: int
    status: str
    baseline_name: str | None
    candidate_name: str | None
    baseline_path: Path | None
    candidate_path: Path | None
    deltas: dict[str, object]
    baseline_status: str | None = None
    candidate_status: str | None = None
    issues: tuple[PresentationVisualIssue, ...] = ()
    recommendations: tuple[str, ...] = ()

    def to_payload(self) -> dict[str, object]:
        return {
            "index": self.index,
            "status": self.status,
            "baseline_name": self.baseline_name,
            "candidate_name": self.candidate_name,
            "baseline_path": str(self.baseline_path) if self.baseline_path else None,
            "candidate_path": str(self.candidate_path) if self.candidate_path else None,
            "baseline_status": self.baseline_status,
            "candidate_status": self.candidate_status,
            "deltas": dict(self.deltas),
            "issues": [issue.to_payload() for issue in self.issues],
            "recommendations": list(self.recommendations),
        }


@dataclass(frozen=True)
class PresentationVisualDiffResult:
    status: str
    baseline_path: Path
    candidate_path: Path
    baseline_checked_paths: tuple[Path, ...]
    candidate_checked_paths: tuple[Path, ...]
    pages: tuple[PresentationVisualDiffPage, ...]
    summary: dict[str, object]
    recommendations: tuple[str, ...] = ()
    truncated: bool = False

    def to_payload(self) -> dict[str, object]:
        return {
            "status": self.status,
            "baseline_path": str(self.baseline_path),
            "candidate_path": str(self.candidate_path),
            "baseline_checked_paths": [str(path) for path in self.baseline_checked_paths],
            "candidate_checked_paths": [str(path) for path in self.candidate_checked_paths],
            "compared_count": len(self.pages),
            "truncated": self.truncated,
            "summary": dict(self.summary),
            "pages": [page.to_payload() for page in self.pages],
            "recommendations": list(self.recommendations),
        }


@dataclass(frozen=True)
class PresentationVisualDiffReportResult:
    status: str
    report_path: Path
    manifest_path: Path
    baseline_path: Path
    candidate_path: Path
    summary: dict[str, object]
    recommendations: tuple[str, ...] = ()
    warnings: tuple[str, ...] = ()

    def to_payload(self) -> dict[str, object]:
        return {
            "status": self.status,
            "report_path": str(self.report_path),
            "manifest_path": str(self.manifest_path),
            "baseline_path": str(self.baseline_path),
            "candidate_path": str(self.candidate_path),
            "summary": dict(self.summary),
            "recommendations": list(self.recommendations),
            "warnings": list(self.warnings),
        }


@dataclass(frozen=True)
class PresentationAnnotation:
    page_index: int
    page_name: str
    source: str
    target: str
    severity: str
    code: str
    message: str
    region: dict[str, object]
    recommendation: str | None = None
    path: Path | None = None

    def to_payload(self) -> dict[str, object]:
        payload: dict[str, object] = {
            "page_index": self.page_index,
            "page_name": self.page_name,
            "source": self.source,
            "target": self.target,
            "severity": self.severity,
            "code": self.code,
            "message": self.message,
            "region": dict(self.region),
        }
        if self.recommendation:
            payload["recommendation"] = self.recommendation
        if self.path is not None:
            payload["path"] = str(self.path)
        return payload


@dataclass(frozen=True)
class PresentationAnnotationSpecResult:
    status: str
    mode: str
    annotations: tuple[PresentationAnnotation, ...]
    summary: dict[str, object]
    recommendations: tuple[str, ...] = ()
    target_path: Path | None = None
    baseline_path: Path | None = None
    candidate_path: Path | None = None
    warnings: tuple[str, ...] = ()

    def to_payload(self) -> dict[str, object]:
        payload: dict[str, object] = {
            "status": self.status,
            "mode": self.mode,
            "annotation_count": len(self.annotations),
            "summary": dict(self.summary),
            "annotations": [annotation.to_payload() for annotation in self.annotations],
            "recommendations": list(self.recommendations),
            "warnings": list(self.warnings),
        }
        if self.target_path is not None:
            payload["target_path"] = str(self.target_path)
        if self.baseline_path is not None:
            payload["baseline_path"] = str(self.baseline_path)
        if self.candidate_path is not None:
            payload["candidate_path"] = str(self.candidate_path)
        return payload


@dataclass(frozen=True)
class PresentationAnnotationReportResult:
    status: str
    report_path: Path
    manifest_path: Path
    mode: str
    summary: dict[str, object]
    recommendations: tuple[str, ...] = ()
    warnings: tuple[str, ...] = ()

    def to_payload(self) -> dict[str, object]:
        return {
            "status": self.status,
            "report_path": str(self.report_path),
            "manifest_path": str(self.manifest_path),
            "mode": self.mode,
            "summary": dict(self.summary),
            "recommendations": list(self.recommendations),
            "warnings": list(self.warnings),
        }


def build_presentation_plan(
    *,
    content: str,
    title: str | None = None,
    audience: str | None = None,
    style: str | None = None,
    max_slides: int = 12,
) -> PresentationPlan:
    bounded_max_slides = max(1, min(max_slides, MAX_SLIDES))
    sections = _markdown_sections(content)
    warnings: list[str] = []
    if not sections:
        sections = [{"title": title or "Untitled", "body": _plain_chunks(content)}]
        warnings.append("no markdown headings found; generated a simple outline from paragraphs")
    deck_title = (title or sections[0]["title"] or "Untitled").strip()
    slides: list[dict[str, object]] = []
    for index, section in enumerate(sections[:bounded_max_slides], start=1):
        raw_body = list(section.get("body") or [])
        bullets = _body_to_bullets(raw_body)
        slide = {
            "index": index,
            "title": str(section.get("title") or f"Slide {index}")[:140],
            "bullets": bullets[:MAX_BULLETS_PER_SLIDE],
            "bullet_count": len(bullets),
            "truncated": len(bullets) > MAX_BULLETS_PER_SLIDE,
        }
        if audience:
            slide["audience"] = audience.strip()[:120]
        if style:
            slide["style"] = style.strip()[:120]
        slides.append(slide)
    if len(sections) > bounded_max_slides:
        warnings.append(f"truncated {len(sections) - bounded_max_slides} section(s) after max_slides")
    return PresentationPlan(
        title=deck_title,
        slides=tuple(slides),
        source="markdown_outline",
        warnings=tuple(warnings),
    )


def build_topic_presentation_content(
    *,
    title: str | None = None,
    audience: str | None = None,
    style: str | None = None,
    max_slides: int = 8,
) -> str:
    deck_title = (title or "Presentation").strip() or "Presentation"
    deck_audience = (audience or "general stakeholders").strip() or "general stakeholders"
    deck_style = (style or "clean editorial").strip() or "clean editorial"
    bounded_slides = max(4, min(max_slides, 10))
    sections = [
        (
            "Executive framing",
            [
                f"Explain why {deck_title} matters to {deck_audience}.",
                "State the operating problem, current friction, and desired business or technical outcome.",
                f"Use a {deck_style} visual system with a clear claim and proof object on every slide.",
            ],
        ),
        (
            "Core idea",
            [
                f"Define {deck_title} in one concrete sentence.",
                "Separate user-facing value from underlying platform mechanisms.",
                "Show the smallest end-to-end workflow that makes the concept tangible.",
            ],
        ),
        (
            "System architecture",
            [
                "Map inputs, context assembly, policy gates, tool execution, memory, and final response boundaries.",
                "Call out synchronous path versus background maintenance so delivery latency is explicit.",
                "Represent external systems as governed adapters rather than hidden side effects.",
            ],
        ),
        (
            "Capability model",
            [
                "Group capabilities into tools, memory, skills, documents, browser, MCP, and automation.",
                "Show how discovery, validation, and runtime visibility protect the main agent loop.",
                "Tie each capability to a measurable user outcome.",
            ],
        ),
        (
            "Operational workflow",
            [
                "Follow request intake, planning, execution, artifact creation, verification, and handoff.",
                "Keep scratch files inside thread-scoped workspace and only final deliverables in the requested directory.",
                "Make failure recovery visible with typed errors and actionable next steps.",
            ],
        ),
        (
            "Governance and safety",
            [
                "Use path boundaries, approval policy, tool schemas, and preflight checks as runtime contracts.",
                "Record evidence for generated files, skill usage, and verification results.",
                "Keep privileged or high-risk actions narrow, auditable, and reversible.",
            ],
        ),
        (
            "Quality bar",
            [
                "Pass a contact-sheet test: distinct slide rhythms, readable hierarchy, and coherent visual language.",
                "Prefer diagrams, metrics, and visual proof over plain bullet pages.",
                "Validate the exported deck for slide count, titles, visual density, and embedded assets.",
            ],
        ),
        (
            "Next steps",
            [
                "Ship the first complete deck, inspect the rendered output, and iterate weakest slides.",
                "Promote reusable workflow lessons into governed skills or procedures.",
                "Track performance, reliability, and user-facing quality as platform metrics.",
            ],
        ),
    ][:bounded_slides]
    lines = [
        f"# {deck_title}",
        "",
        f"- Audience: {deck_audience}",
        f"- Visual style: {deck_style}",
        "- Delivery target: polished editable PowerPoint with diagrams, visual rhythm, and no user-directory scratch artifacts.",
        "",
    ]
    for heading, bullets in sections:
        lines.extend([f"## {heading}", ""])
        lines.extend(f"- {bullet}" for bullet in bullets)
        lines.append("")
    return "\n".join(lines).strip() + "\n"


def build_presentation_contract(
    *,
    content: str,
    title: str | None = None,
    audience: str | None = None,
    style: str | None = None,
    max_slides: int = 12,
    format: str = "ppt169",
    primary_color: str | None = None,
    accent_color: str | None = None,
    font_family: str | None = None,
) -> PresentationContract:
    plan = build_presentation_plan(
        content=content,
        title=title,
        audience=audience,
        style=style,
        max_slides=max_slides,
    )
    canvas = _canvas_for_format(format)
    theme = {
        "primary": _normalize_hex_color(primary_color, fallback="#3A506B"),
        "accent": _normalize_hex_color(accent_color, fallback="#5BC0BE"),
        "background": "#CDEDF6",
        "text": "#1F2A37",
        "mode": style.strip()[:120] if style else "clean editorial",
    }
    typography = {
        "font_family": (font_family or '"Microsoft YaHei", Arial, sans-serif').strip(),
        "title_size": 34,
        "body_size": 20,
        "caption_size": 13,
    }
    slides: list[dict[str, object]] = []
    for slide in plan.slides:
        index = int(slide.get("index") or len(slides) + 1)
        bullet_count = int(slide.get("bullet_count") or len(slide.get("bullets") or []))
        rhythm = "anchor" if index == 1 else "breathing" if bullet_count <= 2 else "dense"
        slide_contract = {
            "index": index,
            "title": slide.get("title") or f"Slide {index}",
            "bullets": list(slide.get("bullets") or []),
            "rhythm": rhythm,
            "layout": _layout_for_slide(index=index, rhythm=rhythm, bullet_count=bullet_count),
            "chart": None,
            "image": None,
        }
        if audience:
            slide_contract["audience"] = audience.strip()[:120]
        slides.append(slide_contract)
    warnings = list(plan.warnings)
    if not slides:
        warnings.append("no slides generated")
    return PresentationContract(
        title=plan.title,
        canvas=canvas,
        theme=theme,
        typography=typography,
        slides=tuple(slides),
        warnings=tuple(warnings),
    )


def preflight_presentation_artifact(path: Path, *, contract: PresentationContract | None = None) -> PresentationPreflightResult:
    suffix = path.suffix.lower()
    if suffix == ".pptx":
        return _preflight_pptx(path, contract=contract)
    if suffix == ".svg":
        return _preflight_svg(path, contract=contract)
    raise ValueError(f"unsupported presentation preflight format: {suffix or path.name}")


def inspect_presentation_template(template_dir: Path) -> PresentationTemplateInfo:
    if not template_dir.exists() or not template_dir.is_dir():
        raise ValueError(f"presentation template directory does not exist: {template_dir}")
    warnings: list[str] = []
    design_spec = template_dir / "design_spec.md"
    if not design_spec.exists():
        warnings.append("missing design_spec.md")
    svg_paths = sorted(template_dir.glob("*.svg"))
    if not svg_paths:
        warnings.append("no top-level SVG pages found")
    pages: list[dict[str, object]] = []
    viewboxes: set[str] = set()
    for svg_path in svg_paths:
        page: dict[str, object] = {
            "name": svg_path.stem,
            "path": str(svg_path),
        }
        try:
            root = ET.fromstring(svg_path.read_text(encoding="utf-8"))
            viewbox = root.attrib.get("viewBox") or root.attrib.get("viewbox")
            if viewbox:
                page["viewBox"] = viewbox
                viewboxes.add(viewbox)
        except Exception as exc:  # noqa: BLE001
            page["error"] = f"invalid svg: {exc}"
            warnings.append(f"{svg_path.name}: invalid SVG")
        pages.append(page)
    asset_paths = sorted(
        path
        for pattern in ("*.png", "*.jpg", "*.jpeg", "*.webp")
        for path in template_dir.glob(pattern)
    )
    assets = tuple(
        {
            "name": path.name,
            "path": str(path),
            "extension": path.suffix.lower(),
            "size_bytes": path.stat().st_size if path.exists() else 0,
        }
        for path in asset_paths
    )
    metadata = _read_template_metadata(design_spec if design_spec.exists() else None)
    if viewboxes:
        metadata["viewBoxes"] = sorted(viewboxes)
    valid = design_spec.exists() and bool(svg_paths) and not any("error" in page for page in pages)
    return PresentationTemplateInfo(
        template_id=template_dir.name,
        template_path=template_dir,
        valid=valid,
        design_spec_path=design_spec if design_spec.exists() else None,
        pages=tuple(pages),
        assets=assets,
        metadata=metadata,
        warnings=tuple(warnings),
    )


def import_presentation_template(
    *,
    template_dir: Path,
    project_dir: Path,
    overwrite: bool = False,
    include_assets: bool = True,
) -> PresentationTemplateImportResult:
    info = inspect_presentation_template(template_dir)
    if not info.valid:
        raise ValueError(f"invalid presentation template: {'; '.join(info.warnings) or template_dir}")
    project_templates = project_dir / "templates"
    project_images = project_dir / "images"
    project_templates.mkdir(parents=True, exist_ok=True)
    copied: list[Path] = []
    skipped: list[Path] = []
    source_paths = [Path(str(info.design_spec_path))] if info.design_spec_path is not None else []
    source_paths.extend(Path(str(page["path"])) for page in info.pages)
    for source in source_paths:
        target = project_templates / source.name
        if target.exists() and not overwrite:
            skipped.append(target)
            continue
        shutil.copy2(source, target)
        copied.append(target)
    if include_assets:
        for asset in info.assets:
            source = Path(str(asset["path"]))
            target = project_images / source.name
            if target.exists() and not overwrite:
                skipped.append(target)
                continue
            project_images.mkdir(parents=True, exist_ok=True)
            shutil.copy2(source, target)
            copied.append(target)
    return PresentationTemplateImportResult(
        status="imported" if copied else "skipped",
        template_id=info.template_id,
        project_path=project_dir,
        template_dir=project_templates,
        copied_paths=tuple(copied),
        skipped_paths=tuple(skipped),
        warnings=info.warnings,
    )


def build_presentation_template_preview(
    *,
    template_dir: Path,
    output_dir: Path,
    max_pages: int = 12,
    include_preflight: bool = True,
) -> PresentationTemplatePreviewResult:
    info = inspect_presentation_template(template_dir)
    if info.design_spec_path is None:
        raise ValueError("presentation template preview requires design_spec.md")
    if not info.pages:
        raise ValueError("presentation template preview requires at least one top-level SVG page")

    bounded_max_pages = max(1, min(max_pages, MAX_TEMPLATE_PREVIEW_PAGES))
    preview_root = output_dir / "presentation-template-previews" / _safe_path_segment(info.template_id)
    preview_root.mkdir(parents=True, exist_ok=True)
    manifest_path = preview_root / "manifest.json"
    preview_path = preview_root / "index.html"

    warnings = list(info.warnings)
    manifest_pages: list[dict[str, object]] = []
    html_pages: list[dict[str, object]] = []
    selected_pages = list(info.pages[:bounded_max_pages])
    if len(info.pages) > bounded_max_pages:
        warnings.append(f"preview truncated {len(info.pages) - bounded_max_pages} page(s) after max_pages")

    for index, page in enumerate(selected_pages, start=1):
        source = Path(str(page["path"]))
        page_warnings: list[str] = []
        data_uri: str | None = None
        size_bytes = source.stat().st_size if source.exists() else 0
        if not source.exists():
            page_warnings.append("source SVG is missing")
        elif size_bytes > MAX_TEMPLATE_PREVIEW_SVG_BYTES:
            page_warnings.append("source SVG is too large for inline preview")
        else:
            data_uri = _svg_data_uri(source)

        preflight_payload: dict[str, object] | None = None
        if include_preflight and source.exists():
            preflight_payload = preflight_presentation_artifact(source).to_payload()

        page_payload: dict[str, object] = {
            "index": index,
            "name": str(page.get("name") or source.stem),
            "source_path": str(source),
            "viewBox": page.get("viewBox"),
            "size_bytes": size_bytes,
            "warnings": page_warnings,
        }
        if preflight_payload is not None:
            page_payload["preflight"] = _compact_preflight_payload(preflight_payload)
        manifest_pages.append(page_payload)
        html_pages.append({**page_payload, "preview_data_uri": data_uri})

    status = _template_preview_status(manifest_pages, warnings)
    manifest = {
        "status": status,
        "template_id": info.template_id,
        "template_path": str(info.template_path),
        "design_spec_path": str(info.design_spec_path),
        "metadata": dict(info.metadata),
        "pages": manifest_pages,
        "warnings": warnings,
    }
    manifest_path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")
    preview_path.write_text(
        _render_template_preview_html(
            template_id=info.template_id,
            metadata=info.metadata,
            pages=html_pages,
            status=status,
            warnings=warnings,
        ),
        encoding="utf-8",
    )
    return PresentationTemplatePreviewResult(
        status=status,
        template_id=info.template_id,
        preview_path=preview_path,
        manifest_path=manifest_path,
        pages=tuple(manifest_pages),
        warnings=tuple(warnings),
    )


def review_presentation_visuals(
    *,
    target_path: Path,
    contract_payload: dict[str, object] | None = None,
    max_pages: int = 20,
    include_preflight: bool = True,
) -> PresentationVisualReviewResult:
    svg_paths = _visual_review_svg_paths(target_path, max_pages=max_pages)
    if not svg_paths:
        raise ValueError(f"presentation visual review found no SVG pages: {target_path}")

    contract = _normalize_contract_payload(contract_payload)
    pages: list[PresentationVisualPageReview] = []
    for index, svg_path in enumerate(svg_paths, start=1):
        pages.append(
            _review_svg_page(
                path=svg_path,
                index=index,
                contract=contract,
                include_preflight=include_preflight,
            )
        )

    error_count = sum(1 for page in pages for issue in page.issues if issue.severity == "error")
    warning_count = sum(1 for page in pages for issue in page.issues if issue.severity == "warning")
    page_status_counts = {
        "passed": sum(1 for page in pages if page.status == "passed"),
        "warning": sum(1 for page in pages if page.status == "warning"),
        "failed": sum(1 for page in pages if page.status == "failed"),
    }
    palette: set[str] = set()
    total_text_chars = 0
    total_text_nodes = 0
    for page in pages:
        colors = page.metrics.get("colors")
        if isinstance(colors, list):
            palette.update(str(color) for color in colors)
        total_text_chars += int(page.metrics.get("text_characters") or 0)
        total_text_nodes += int(page.metrics.get("text_node_count") or 0)
    status = "failed" if error_count else "warning" if warning_count else "passed"
    recommendations = _visual_review_recommendations(pages, error_count=error_count, warning_count=warning_count)
    summary = {
        "error_count": error_count,
        "warning_count": warning_count,
        "page_status_counts": page_status_counts,
        "palette": sorted(palette)[:24],
        "palette_truncated": len(palette) > 24,
        "total_text_characters": total_text_chars,
        "total_text_nodes": total_text_nodes,
        "contract_slide_count": len(contract.get("slides") or []),
    }
    all_svg_paths = _list_visual_review_svg_paths(target_path)
    return PresentationVisualReviewResult(
        status=status,
        target_path=target_path,
        checked_paths=tuple(svg_paths),
        pages=tuple(pages),
        summary=summary,
        recommendations=tuple(recommendations),
        truncated=len(all_svg_paths) > len(svg_paths),
    )


def build_presentation_review_report(
    *,
    target_path: Path,
    output_dir: Path,
    contract_payload: dict[str, object] | None = None,
    max_pages: int = 20,
    include_preflight: bool = True,
    include_svg_previews: bool = True,
    report_id: str | None = None,
) -> PresentationReviewReportResult:
    review = review_presentation_visuals(
        target_path=target_path,
        contract_payload=contract_payload,
        max_pages=max_pages,
        include_preflight=include_preflight,
    )
    report_slug = _safe_path_segment(report_id or f"{target_path.stem or target_path.name}-visual-review")
    report_root = output_dir / "presentation-review-reports" / report_slug
    report_root.mkdir(parents=True, exist_ok=True)
    report_path = report_root / "index.html"
    manifest_path = report_root / "manifest.json"

    manifest_pages: list[dict[str, object]] = []
    html_pages: list[dict[str, object]] = []
    warnings: list[str] = []
    for page in review.pages:
        display_path = _artifact_display_path(page.path, root=target_path)
        size_bytes = page.path.stat().st_size if page.path.exists() else 0
        preview_data_uri: str | None = None
        page_warnings: list[str] = []
        if include_svg_previews:
            if not page.path.exists():
                page_warnings.append("source SVG is missing")
            elif size_bytes > MAX_TEMPLATE_PREVIEW_SVG_BYTES:
                page_warnings.append("source SVG is too large for inline report preview")
            else:
                preview_data_uri = _svg_data_uri(page.path)
        page_payload = {
            "index": page.index,
            "name": page.name,
            "path": display_path,
            "status": page.status,
            "viewBox": page.view_box,
            "metrics": dict(page.metrics),
            "contract": dict(page.contract),
            "issues": [issue.to_payload() for issue in page.issues],
            "recommendations": list(page.recommendations),
            "preflight": _compact_preflight_payload(page.preflight) if isinstance(page.preflight, dict) else None,
            "size_bytes": size_bytes,
            "warnings": page_warnings,
        }
        manifest_pages.append(page_payload)
        html_pages.append({**page_payload, "preview_data_uri": preview_data_uri})
        warnings.extend(page_warnings)

    manifest = {
        "status": review.status,
        "target": _artifact_display_path(target_path, root=target_path.parent if target_path.parent else target_path),
        "checked_count": len(review.checked_paths),
        "truncated": review.truncated,
        "summary": dict(review.summary),
        "recommendations": list(review.recommendations),
        "pages": manifest_pages,
        "warnings": warnings,
    }
    manifest_path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")
    report_path.write_text(
        _render_visual_review_report_html(
            title=report_slug,
            status=review.status,
            summary=review.summary,
            recommendations=list(review.recommendations),
            pages=html_pages,
            warnings=warnings,
        ),
        encoding="utf-8",
    )
    return PresentationReviewReportResult(
        status=review.status,
        report_path=report_path,
        manifest_path=manifest_path,
        target_path=target_path,
        checked_paths=review.checked_paths,
        summary=review.summary,
        recommendations=review.recommendations,
        warnings=tuple(warnings),
    )


def diff_presentation_visuals(
    *,
    baseline_path: Path,
    candidate_path: Path,
    contract_payload: dict[str, object] | None = None,
    max_pages: int = 20,
    include_preflight: bool = True,
) -> PresentationVisualDiffResult:
    baseline = review_presentation_visuals(
        target_path=baseline_path,
        contract_payload=contract_payload,
        max_pages=max_pages,
        include_preflight=include_preflight,
    )
    candidate = review_presentation_visuals(
        target_path=candidate_path,
        contract_payload=contract_payload,
        max_pages=max_pages,
        include_preflight=include_preflight,
    )
    pages: list[PresentationVisualDiffPage] = []
    max_count = max(len(baseline.pages), len(candidate.pages))
    for offset in range(max_count):
        baseline_page = baseline.pages[offset] if offset < len(baseline.pages) else None
        candidate_page = candidate.pages[offset] if offset < len(candidate.pages) else None
        pages.append(_diff_visual_page(index=offset + 1, baseline_page=baseline_page, candidate_page=candidate_page))

    error_count = sum(1 for page in pages for issue in page.issues if issue.severity == "error")
    warning_count = sum(1 for page in pages for issue in page.issues if issue.severity == "warning")
    changed_pages = sum(1 for page in pages if page.status != "unchanged")
    status = "failed" if error_count else "warning" if warning_count else "passed"
    summary = {
        "error_count": error_count,
        "warning_count": warning_count,
        "changed_pages": changed_pages,
        "baseline_page_count": len(baseline.pages),
        "candidate_page_count": len(candidate.pages),
        "baseline_status": baseline.status,
        "candidate_status": candidate.status,
    }
    return PresentationVisualDiffResult(
        status=status,
        baseline_path=baseline_path,
        candidate_path=candidate_path,
        baseline_checked_paths=baseline.checked_paths,
        candidate_checked_paths=candidate.checked_paths,
        pages=tuple(pages),
        summary=summary,
        recommendations=tuple(_visual_diff_recommendations(pages, error_count=error_count, warning_count=warning_count)),
        truncated=baseline.truncated or candidate.truncated,
    )


def build_presentation_visual_diff_report(
    *,
    baseline_path: Path,
    candidate_path: Path,
    output_dir: Path,
    contract_payload: dict[str, object] | None = None,
    max_pages: int = 20,
    include_preflight: bool = True,
    include_svg_previews: bool = True,
    report_id: str | None = None,
) -> PresentationVisualDiffReportResult:
    diff = diff_presentation_visuals(
        baseline_path=baseline_path,
        candidate_path=candidate_path,
        contract_payload=contract_payload,
        max_pages=max_pages,
        include_preflight=include_preflight,
    )
    report_slug = _safe_path_segment(report_id or f"{baseline_path.stem or baseline_path.name}-to-{candidate_path.stem or candidate_path.name}-visual-diff")
    report_root = output_dir / "presentation-visual-diffs" / report_slug
    report_root.mkdir(parents=True, exist_ok=True)
    report_path = report_root / "index.html"
    manifest_path = report_root / "manifest.json"
    warnings: list[str] = []
    manifest_pages: list[dict[str, object]] = []
    html_pages: list[dict[str, object]] = []
    for page in diff.pages:
        page_payload = _visual_diff_page_manifest(page, baseline_root=baseline_path, candidate_root=candidate_path)
        baseline_preview = None
        candidate_preview = None
        if include_svg_previews:
            baseline_preview, baseline_warnings = _safe_svg_preview_data_uri(page.baseline_path, label="baseline")
            candidate_preview, candidate_warnings = _safe_svg_preview_data_uri(page.candidate_path, label="candidate")
            warnings.extend(baseline_warnings)
            warnings.extend(candidate_warnings)
            page_payload["warnings"] = baseline_warnings + candidate_warnings
        manifest_pages.append(page_payload)
        html_pages.append(
            {
                **page_payload,
                "baseline_preview_data_uri": baseline_preview,
                "candidate_preview_data_uri": candidate_preview,
            }
        )

    manifest = {
        "status": diff.status,
        "summary": dict(diff.summary),
        "baseline": _artifact_display_path(baseline_path, root=baseline_path.parent if baseline_path.parent else baseline_path),
        "candidate": _artifact_display_path(candidate_path, root=candidate_path.parent if candidate_path.parent else candidate_path),
        "truncated": diff.truncated,
        "recommendations": list(diff.recommendations),
        "pages": manifest_pages,
        "warnings": warnings,
    }
    manifest_path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")
    report_path.write_text(
        _render_visual_diff_report_html(
            title=report_slug,
            status=diff.status,
            summary=diff.summary,
            recommendations=list(diff.recommendations),
            pages=html_pages,
            warnings=warnings,
        ),
        encoding="utf-8",
    )
    return PresentationVisualDiffReportResult(
        status=diff.status,
        report_path=report_path,
        manifest_path=manifest_path,
        baseline_path=baseline_path,
        candidate_path=candidate_path,
        summary=diff.summary,
        recommendations=diff.recommendations,
        warnings=tuple(warnings),
    )


def build_presentation_annotation_spec(
    *,
    target_path: Path | None = None,
    baseline_path: Path | None = None,
    candidate_path: Path | None = None,
    contract_payload: dict[str, object] | None = None,
    max_pages: int = 20,
    include_preflight: bool = True,
    mode: str = "auto",
) -> PresentationAnnotationSpecResult:
    normalized_mode = str(mode or "auto").strip().lower()
    if normalized_mode not in {"auto", "review", "diff"}:
        raise ValueError("mode must be auto, review, or diff")
    use_diff = normalized_mode == "diff" or (normalized_mode == "auto" and baseline_path is not None and candidate_path is not None)
    if use_diff:
        if baseline_path is None or candidate_path is None:
            raise ValueError("diff annotation mode requires baseline_path and candidate_path")
        diff = diff_presentation_visuals(
            baseline_path=baseline_path,
            candidate_path=candidate_path,
            contract_payload=contract_payload,
            max_pages=max_pages,
            include_preflight=include_preflight,
        )
        annotations = _annotations_from_visual_diff(diff)
        summary = {
            **dict(diff.summary),
            "annotation_count": len(annotations),
        }
        return PresentationAnnotationSpecResult(
            status=diff.status,
            mode="diff",
            annotations=tuple(annotations),
            summary=summary,
            recommendations=diff.recommendations,
            baseline_path=baseline_path,
            candidate_path=candidate_path,
            warnings=(),
        )
    if target_path is None:
        raise ValueError("review annotation mode requires target_path")
    review = review_presentation_visuals(
        target_path=target_path,
        contract_payload=contract_payload,
        max_pages=max_pages,
        include_preflight=include_preflight,
    )
    annotations = _annotations_from_visual_review(review)
    summary = {
        **dict(review.summary),
        "annotation_count": len(annotations),
    }
    return PresentationAnnotationSpecResult(
        status=review.status,
        mode="review",
        annotations=tuple(annotations),
        summary=summary,
        recommendations=review.recommendations,
        target_path=target_path,
        warnings=(),
    )


def build_presentation_annotation_report(
    *,
    output_dir: Path,
    target_path: Path | None = None,
    baseline_path: Path | None = None,
    candidate_path: Path | None = None,
    contract_payload: dict[str, object] | None = None,
    max_pages: int = 20,
    include_preflight: bool = True,
    include_svg_previews: bool = True,
    report_id: str | None = None,
    mode: str = "auto",
) -> PresentationAnnotationReportResult:
    spec = build_presentation_annotation_spec(
        target_path=target_path,
        baseline_path=baseline_path,
        candidate_path=candidate_path,
        contract_payload=contract_payload,
        max_pages=max_pages,
        include_preflight=include_preflight,
        mode=mode,
    )
    base_name = report_id or f"{spec.mode}-annotations"
    report_slug = _safe_path_segment(base_name)
    report_root = output_dir / "presentation-annotations" / report_slug
    report_root.mkdir(parents=True, exist_ok=True)
    report_path = report_root / "index.html"
    manifest_path = report_root / "manifest.json"
    warnings: list[str] = []
    pages = _annotation_pages_for_report(spec)
    html_pages: list[dict[str, object]] = []
    for page in pages:
        source_path = page.get("path")
        preview_data_uri = None
        if include_svg_previews and isinstance(source_path, Path):
            preview_data_uri, preview_warnings = _safe_svg_preview_data_uri(source_path, label=str(page.get("name") or "page"))
            warnings.extend(preview_warnings)
        html_pages.append({**page, "preview_data_uri": preview_data_uri})
    manifest = {
        **spec.to_payload(),
        "annotations": [_annotation_manifest_payload(annotation, spec=spec) for annotation in spec.annotations],
        "pages": [_annotation_page_manifest_payload(page) for page in pages],
        "warnings": list(spec.warnings) + warnings,
    }
    manifest_path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")
    report_path.write_text(
        _render_annotation_report_html(
            title=report_slug,
            status=spec.status,
            mode=spec.mode,
            summary=spec.summary,
            recommendations=list(spec.recommendations),
            pages=html_pages,
            annotations=list(spec.annotations),
            warnings=list(spec.warnings) + warnings,
        ),
        encoding="utf-8",
    )
    return PresentationAnnotationReportResult(
        status=spec.status,
        report_path=report_path,
        manifest_path=manifest_path,
        mode=spec.mode,
        summary=spec.summary,
        recommendations=spec.recommendations,
        warnings=tuple(list(spec.warnings) + warnings),
    )


def _markdown_sections(content: str) -> list[dict[str, object]]:
    sections: list[dict[str, object]] = []
    current: dict[str, object] | None = None
    for line in content.splitlines():
        heading = re.match(r"^(#{1,3})\s+(.+?)\s*$", line.strip())
        if heading:
            if current is not None:
                sections.append(current)
            current = {"title": heading.group(2).strip(), "body": []}
            continue
        if current is None:
            if not line.strip():
                continue
            current = {"title": "Overview", "body": []}
        current_body = current.setdefault("body", [])
        if isinstance(current_body, list):
            current_body.append(line)
    if current is not None:
        sections.append(current)
    return [
        {
            "title": str(section.get("title") or "Untitled").strip(),
            "body": [str(line) for line in section.get("body") or []],
        }
        for section in sections
        if str(section.get("title") or "").strip() or any(str(line).strip() for line in section.get("body") or [])
    ]


def _visual_review_svg_paths(target_path: Path, *, max_pages: int) -> list[Path]:
    bounded_max_pages = max(1, min(max_pages, MAX_VISUAL_REVIEW_PAGES))
    return _list_visual_review_svg_paths(target_path)[:bounded_max_pages]


def _list_visual_review_svg_paths(target_path: Path) -> list[Path]:
    if target_path.is_file():
        return [target_path] if target_path.suffix.lower() == ".svg" else []
    if not target_path.exists() or not target_path.is_dir():
        return []
    svg_output = target_path / "svg_output"
    search_dir = svg_output if svg_output.exists() and svg_output.is_dir() else target_path
    return sorted(search_dir.glob("*.svg"), key=lambda item: item.name)


def _normalize_contract_payload(payload: dict[str, object] | None) -> dict[str, object]:
    if not isinstance(payload, dict):
        return {}
    slides = payload.get("slides")
    normalized = {
        "canvas": dict(payload.get("canvas") or {}) if isinstance(payload.get("canvas"), dict) else {},
        "theme": dict(payload.get("theme") or {}) if isinstance(payload.get("theme"), dict) else {},
        "typography": dict(payload.get("typography") or {}) if isinstance(payload.get("typography"), dict) else {},
        "slides": list(slides) if isinstance(slides, list) else [],
    }
    return normalized


def _review_svg_page(
    *,
    path: Path,
    index: int,
    contract: dict[str, object],
    include_preflight: bool,
) -> PresentationVisualPageReview:
    issues: list[PresentationVisualIssue] = []
    recommendations: list[str] = []
    preflight_payload = preflight_presentation_artifact(path).to_payload() if include_preflight else None
    if preflight_payload is not None:
        for issue in preflight_payload.get("issues") or []:
            if isinstance(issue, dict):
                issues.append(
                    PresentationVisualIssue(
                        severity=str(issue.get("severity") or "warning"),
                        code=f"preflight_{issue.get('code') or 'issue'}",
                        message=str(issue.get("message") or issue.get("code") or "preflight issue"),
                    )
                )
    try:
        content = path.read_text(encoding="utf-8")
        root = ET.fromstring(content)
    except Exception as exc:  # noqa: BLE001
        issues.append(PresentationVisualIssue("error", "invalid_svg", f"unable to parse SVG: {exc}"))
        return PresentationVisualPageReview(
            index=index,
            name=path.stem,
            path=path,
            status="failed",
            view_box=None,
            metrics={},
            contract=_page_contract_payload(contract, index=index),
            issues=tuple(issues),
            recommendations=("Regenerate or repair the SVG before export.",),
            preflight=preflight_payload,
        )

    metrics = _svg_visual_metrics(content, root=root)
    page_contract = _page_contract_payload(contract, index=index)
    issues.extend(_visual_contract_issues(metrics=metrics, contract=contract, page_contract=page_contract))
    issues.extend(_visual_density_issues(metrics=metrics, page_contract=page_contract))
    recommendations.extend(_page_visual_recommendations(metrics=metrics, issues=issues, page_contract=page_contract))
    status = "failed" if any(issue.severity == "error" for issue in issues) else "warning" if issues else "passed"
    return PresentationVisualPageReview(
        index=index,
        name=path.stem,
        path=path,
        status=status,
        view_box=str(metrics.get("viewBox") or "") or None,
        metrics=metrics,
        contract=page_contract,
        issues=tuple(issues),
        recommendations=tuple(recommendations),
        preflight=preflight_payload,
    )


def _svg_visual_metrics(content: str, *, root: ET.Element) -> dict[str, object]:
    viewbox = root.attrib.get("viewBox") or root.attrib.get("viewbox")
    text_values = [
        " ".join((element.itertext() if hasattr(element, "itertext") else []))
        for element in root.iter()
        if _xml_local_name(element.tag) == "text"
    ]
    text_values = [value.strip() for value in text_values if value.strip()]
    font_sizes = [
        float(match.group(1))
        for match in re.finditer(r"font-size\s*[:=]\s*[\"']?([0-9.]+)", content, re.IGNORECASE)
        if _safe_float(match.group(1)) is not None
    ]
    colors = sorted(set(color.upper() for color in re.findall(r"#[0-9A-Fa-f]{6}", content)))
    element_counts: dict[str, int] = {}
    for element in root.iter():
        name = _xml_local_name(element.tag)
        element_counts[name] = element_counts.get(name, 0) + 1
    image_refs = len(re.findall(r"<image\b", content, re.IGNORECASE))
    group_ids = [
        str(element.attrib.get("id") or "")
        for element in root.iter()
        if _xml_local_name(element.tag) == "g" and element.attrib.get("id")
    ]
    width, height = _viewbox_size(viewbox)
    shape_count = sum(element_counts.get(name, 0) for name in ("rect", "circle", "ellipse", "line", "polyline", "polygon", "path"))
    return {
        "viewBox": viewbox,
        "width": width,
        "height": height,
        "element_count": sum(element_counts.values()),
        "shape_count": shape_count,
        "image_count": image_refs,
        "text_node_count": len(text_values),
        "text_characters": sum(len(value) for value in text_values),
        "max_text_characters": max((len(value) for value in text_values), default=0),
        "font_sizes": sorted(set(font_sizes))[:24],
        "font_size_count": len(set(font_sizes)),
        "colors": colors[:32],
        "color_count": len(colors),
        "group_ids": group_ids[:32],
        "group_id_count": len(group_ids),
        "element_counts": dict(sorted(element_counts.items())[:32]),
        "render_signature": _svg_render_signature(root=root, width=width, height=height),
    }


def _page_contract_payload(contract: dict[str, object], *, index: int) -> dict[str, object]:
    slides = contract.get("slides")
    slide = {}
    if isinstance(slides, list) and index - 1 < len(slides) and isinstance(slides[index - 1], dict):
        slide = dict(slides[index - 1])
    return {
        "expected_viewBox": (contract.get("canvas") or {}).get("viewBox") if isinstance(contract.get("canvas"), dict) else None,
        "theme": dict(contract.get("theme") or {}) if isinstance(contract.get("theme"), dict) else {},
        "typography": dict(contract.get("typography") or {}) if isinstance(contract.get("typography"), dict) else {},
        "slide": {
            "title": slide.get("title"),
            "rhythm": slide.get("rhythm"),
            "layout": slide.get("layout"),
            "bullet_count": len(slide.get("bullets") or []) if isinstance(slide.get("bullets"), list) else None,
        },
    }


def _visual_contract_issues(
    *,
    metrics: dict[str, object],
    contract: dict[str, object],
    page_contract: dict[str, object],
) -> list[PresentationVisualIssue]:
    issues: list[PresentationVisualIssue] = []
    expected_viewbox = page_contract.get("expected_viewBox")
    if expected_viewbox and metrics.get("viewBox") != expected_viewbox:
        issues.append(
            PresentationVisualIssue(
                "warning",
                "viewbox_contract_mismatch",
                f"SVG viewBox '{metrics.get('viewBox')}' does not match contract '{expected_viewbox}'.",
            )
        )
    theme = contract.get("theme") if isinstance(contract.get("theme"), dict) else {}
    expected_colors = {
        str(theme.get(key)).upper()
        for key in ("primary", "accent", "background", "text")
        if isinstance(theme.get(key), str) and re.fullmatch(r"#[0-9A-Fa-f]{6}", str(theme.get(key)))
    }
    colors = {str(color).upper() for color in metrics.get("colors") or []}
    if expected_colors and not (colors & expected_colors):
        issues.append(
            PresentationVisualIssue(
                "warning",
                "palette_contract_drift",
                "SVG does not use any locked theme color from the presentation contract.",
            )
        )
    typography = contract.get("typography") if isinstance(contract.get("typography"), dict) else {}
    title_size = _safe_float(typography.get("title_size"))
    body_size = _safe_float(typography.get("body_size"))
    font_sizes = [float(value) for value in metrics.get("font_sizes") or [] if _safe_float(value) is not None]
    if title_size and body_size and font_sizes:
        minimum = max(body_size * 0.45, 6)
        maximum = max(title_size * 1.8, body_size * 4)
        outside = [size for size in font_sizes if size < minimum or size > maximum]
        if outside:
            issues.append(
                PresentationVisualIssue(
                    "warning",
                    "font_size_outside_contract",
                    f"SVG has font sizes outside the contract ramp: {outside[:8]}.",
                )
            )
    return issues


def _visual_density_issues(
    *,
    metrics: dict[str, object],
    page_contract: dict[str, object],
) -> list[PresentationVisualIssue]:
    issues: list[PresentationVisualIssue] = []
    text_chars = int(metrics.get("text_characters") or 0)
    text_nodes = int(metrics.get("text_node_count") or 0)
    shape_count = int(metrics.get("shape_count") or 0)
    color_count = int(metrics.get("color_count") or 0)
    element_count = int(metrics.get("element_count") or 0)
    slide = page_contract.get("slide") if isinstance(page_contract.get("slide"), dict) else {}
    layout = str(slide.get("layout") or "")
    rhythm = str(slide.get("rhythm") or "")
    if text_chars > 850 or text_nodes > 18:
        issues.append(
            PresentationVisualIssue(
                "warning",
                "text_density_high",
                f"Page has {text_chars} text characters across {text_nodes} text nodes; consider splitting or simplifying.",
            )
        )
    if rhythm in {"anchor", "breathing"} and text_chars > 420:
        issues.append(
            PresentationVisualIssue(
                "warning",
                "rhythm_density_mismatch",
                f"Contract rhythm '{rhythm}' expects a lighter page than the current text density.",
            )
        )
    if layout == "cover" and text_nodes > 8:
        issues.append(PresentationVisualIssue("warning", "cover_too_many_text_nodes", "Cover slide has too many text nodes for a focused first impression."))
    if color_count > 10:
        issues.append(
            PresentationVisualIssue(
                "warning",
                "palette_too_broad",
                f"Page uses {color_count} unique hex colors; consider narrowing the palette.",
            )
        )
    if shape_count > 90 or element_count > 180:
        issues.append(
            PresentationVisualIssue(
                "warning",
                "visual_complexity_high",
                f"Page has {shape_count} shape elements and {element_count} total SVG elements.",
            )
        )
    if shape_count == 0 and text_nodes == 0 and int(metrics.get("image_count") or 0) == 0:
        issues.append(PresentationVisualIssue("error", "page_appears_blank", "SVG has no text, shape, or image content."))
    return issues


def _page_visual_recommendations(
    *,
    metrics: dict[str, object],
    issues: list[PresentationVisualIssue],
    page_contract: dict[str, object],
) -> list[str]:
    recommendations: list[str] = []
    issue_codes = {issue.code for issue in issues}
    if "viewbox_contract_mismatch" in issue_codes:
        recommendations.append("Regenerate the page with the contract canvas viewBox before PPTX conversion.")
    if "palette_contract_drift" in issue_codes or "palette_too_broad" in issue_codes:
        recommendations.append("Re-anchor fills and strokes to the locked primary/accent/background/text colors.")
    if "text_density_high" in issue_codes or "rhythm_density_mismatch" in issue_codes:
        recommendations.append("Split dense content into another slide or reduce bullets to match the declared rhythm.")
    if "visual_complexity_high" in issue_codes:
        recommendations.append("Group or simplify decorative elements before conversion to editable PPTX shapes.")
    if "page_appears_blank" in issue_codes:
        recommendations.append("Re-author the page; do not export a blank SVG.")
    if not recommendations and page_contract.get("slide"):
        recommendations.append("Page structure matches the current lightweight visual review checks.")
    return recommendations[:6]


def _visual_review_recommendations(
    pages: list[PresentationVisualPageReview],
    *,
    error_count: int,
    warning_count: int,
) -> list[str]:
    if error_count:
        return ["Fix failed pages before export; structural errors will likely break PPTX conversion."]
    if warning_count:
        return ["Review warning pages for density, palette drift, and contract mismatch before final delivery."]
    if pages:
        return ["Visual review passed lightweight structure, density, palette, and contract checks."]
    return []


def _diff_visual_page(
    *,
    index: int,
    baseline_page: PresentationVisualPageReview | None,
    candidate_page: PresentationVisualPageReview | None,
) -> PresentationVisualDiffPage:
    issues: list[PresentationVisualIssue] = []
    recommendations: list[str] = []
    if baseline_page is None:
        issues.append(PresentationVisualIssue("warning", "page_added", "Candidate deck has an additional page."))
        return PresentationVisualDiffPage(
            index=index,
            status="added",
            baseline_name=None,
            candidate_name=candidate_page.name if candidate_page else None,
            baseline_path=None,
            candidate_path=candidate_page.path if candidate_page else None,
            baseline_status=None,
            candidate_status=candidate_page.status if candidate_page else None,
            deltas={"page": "added"},
            issues=tuple(issues),
            recommendations=("Confirm the added page is intentional and covered by the deck contract.",),
        )
    if candidate_page is None:
        issues.append(PresentationVisualIssue("error", "page_removed", "Candidate deck is missing a baseline page."))
        return PresentationVisualDiffPage(
            index=index,
            status="removed",
            baseline_name=baseline_page.name,
            candidate_name=None,
            baseline_path=baseline_page.path,
            candidate_path=None,
            baseline_status=baseline_page.status,
            candidate_status=None,
            deltas={"page": "removed"},
            issues=tuple(issues),
            recommendations=("Restore the missing page or update the deck contract before delivery.",),
        )

    deltas = _visual_metric_deltas(baseline_page.metrics, candidate_page.metrics)
    baseline_issue_codes = {issue.code for issue in baseline_page.issues}
    candidate_issue_codes = {issue.code for issue in candidate_page.issues}
    new_issue_codes = sorted(candidate_issue_codes - baseline_issue_codes)
    resolved_issue_codes = sorted(baseline_issue_codes - candidate_issue_codes)
    deltas["new_issue_codes"] = new_issue_codes
    deltas["resolved_issue_codes"] = resolved_issue_codes
    if baseline_page.view_box != candidate_page.view_box:
        issues.append(
            PresentationVisualIssue(
                "warning",
                "viewbox_changed",
                f"Page viewBox changed from '{baseline_page.view_box}' to '{candidate_page.view_box}'.",
            )
        )
    if candidate_page.status == "failed" and baseline_page.status != "failed":
        issues.append(PresentationVisualIssue("error", "candidate_regressed_to_failed", "Candidate page now fails visual review."))
    elif candidate_page.status == "warning" and baseline_page.status == "passed":
        issues.append(PresentationVisualIssue("warning", "candidate_regressed_to_warning", "Candidate page now has visual review warnings."))
    for code in new_issue_codes[:8]:
        issues.append(PresentationVisualIssue("warning", f"new_{code}", f"Candidate introduced issue '{code}'."))
    issues.extend(_visual_metric_regression_issues(deltas))
    recommendations.extend(_visual_diff_page_recommendations(deltas=deltas, issues=issues))
    status = "failed" if any(issue.severity == "error" for issue in issues) else "warning" if issues else "unchanged"
    return PresentationVisualDiffPage(
        index=index,
        status=status,
        baseline_name=baseline_page.name,
        candidate_name=candidate_page.name,
        baseline_path=baseline_page.path,
        candidate_path=candidate_page.path,
        baseline_status=baseline_page.status,
        candidate_status=candidate_page.status,
        deltas=deltas,
        issues=tuple(issues),
        recommendations=tuple(recommendations),
    )


def _visual_metric_deltas(baseline_metrics: dict[str, object], candidate_metrics: dict[str, object]) -> dict[str, object]:
    numeric_fields = [
        "text_characters",
        "text_node_count",
        "shape_count",
        "element_count",
        "image_count",
        "color_count",
        "font_size_count",
        "max_text_characters",
    ]
    deltas: dict[str, object] = {}
    for field in numeric_fields:
        baseline_value = int(baseline_metrics.get(field) or 0)
        candidate_value = int(candidate_metrics.get(field) or 0)
        deltas[field] = {
            "baseline": baseline_value,
            "candidate": candidate_value,
            "delta": candidate_value - baseline_value,
        }
    baseline_colors = {str(color).upper() for color in baseline_metrics.get("colors") or []}
    candidate_colors = {str(color).upper() for color in candidate_metrics.get("colors") or []}
    deltas["colors_added"] = sorted(candidate_colors - baseline_colors)[:16]
    deltas["colors_removed"] = sorted(baseline_colors - candidate_colors)[:16]
    baseline_fonts = {_safe_float(value) for value in baseline_metrics.get("font_sizes") or []}
    candidate_fonts = {_safe_float(value) for value in candidate_metrics.get("font_sizes") or []}
    baseline_fonts.discard(None)
    candidate_fonts.discard(None)
    deltas["font_sizes_added"] = sorted(candidate_fonts - baseline_fonts)[:16]
    deltas["font_sizes_removed"] = sorted(baseline_fonts - candidate_fonts)[:16]
    deltas["render_signature"] = _render_signature_delta(
        baseline_metrics.get("render_signature"),
        candidate_metrics.get("render_signature"),
    )
    return deltas


def _visual_metric_regression_issues(deltas: dict[str, object]) -> list[PresentationVisualIssue]:
    issues: list[PresentationVisualIssue] = []
    text_delta = _delta_value(deltas, "text_characters")
    text_node_delta = _delta_value(deltas, "text_node_count")
    shape_delta = _delta_value(deltas, "shape_count")
    element_delta = _delta_value(deltas, "element_count")
    color_delta = _delta_value(deltas, "color_count")
    if text_delta > 280 or text_node_delta > 8:
        issues.append(PresentationVisualIssue("warning", "text_density_regressed", "Candidate significantly increases text density."))
    if shape_delta > 45 or element_delta > 90:
        issues.append(PresentationVisualIssue("warning", "visual_complexity_regressed", "Candidate significantly increases SVG complexity."))
    if color_delta > 5:
        issues.append(PresentationVisualIssue("warning", "palette_regressed", "Candidate significantly broadens the color palette."))
    if deltas.get("colors_added") and len(deltas.get("colors_added") or []) >= 4:
        issues.append(PresentationVisualIssue("warning", "new_palette_colors", "Candidate introduces several new colors."))
    render_delta = deltas.get("render_signature")
    if isinstance(render_delta, dict):
        occupancy_delta = abs(int(render_delta.get("occupied_cell_delta") or 0))
        changed_ratio = float(render_delta.get("changed_cell_ratio") or 0)
        moved_weight = float(render_delta.get("moved_weight_ratio") or 0)
        color_drift = float(render_delta.get("color_distribution_drift") or 0)
        if changed_ratio >= 0.35 or moved_weight >= 0.30:
            issues.append(PresentationVisualIssue("warning", "render_layout_shift", "Candidate has a large spatial layout shift in the SVG render signature."))
        if occupancy_delta >= 8:
            issues.append(PresentationVisualIssue("warning", "render_occupancy_changed", "Candidate materially changes the occupied visual grid area."))
        if color_drift >= 0.35:
            issues.append(PresentationVisualIssue("warning", "render_color_distribution_shift", "Candidate materially changes rendered color distribution."))
    return issues


def _delta_value(deltas: dict[str, object], field: str) -> int:
    payload = deltas.get(field)
    if not isinstance(payload, dict):
        return 0
    return int(payload.get("delta") or 0)


def _svg_render_signature(*, root: ET.Element, width: float | None, height: float | None) -> dict[str, object]:
    if not width or not height or width <= 0 or height <= 0:
        return {
            "available": False,
            "reason": "missing_viewbox_size",
            "grid": [0] * (VISUAL_SIGNATURE_GRID_COLUMNS * VISUAL_SIGNATURE_GRID_ROWS),
            "occupied_cells": 0,
            "dominant_colors": [],
            "element_count": 0,
        }
    grid = [0.0] * (VISUAL_SIGNATURE_GRID_COLUMNS * VISUAL_SIGNATURE_GRID_ROWS)
    color_weights: dict[str, float] = {}
    counted = 0
    for element in root.iter():
        name = _xml_local_name(element.tag)
        bounds = _svg_element_bounds(element, name=name)
        if bounds is None:
            continue
        x, y, element_width, element_height = bounds
        if element_width <= 0 or element_height <= 0:
            continue
        color = _svg_element_color(element)
        area_weight = max(1.0, min(element_width * element_height, width * height))
        _accumulate_render_grid(grid, x=x, y=y, width=element_width, height=element_height, view_width=width, view_height=height, weight=area_weight)
        color_weights[color] = color_weights.get(color, 0.0) + area_weight
        counted += 1
    total_weight = sum(grid)
    normalized_grid = [round(value / total_weight, 6) if total_weight else 0 for value in grid]
    total_color_weight = sum(color_weights.values())
    dominant_colors = [
        {"color": color, "weight": round(weight / total_color_weight, 6) if total_color_weight else 0}
        for color, weight in sorted(color_weights.items(), key=lambda item: (-item[1], item[0]))[:8]
    ]
    occupied_cells = sum(1 for value in normalized_grid if value > 0)
    return {
        "available": counted > 0,
        "grid_size": [VISUAL_SIGNATURE_GRID_COLUMNS, VISUAL_SIGNATURE_GRID_ROWS],
        "grid": normalized_grid,
        "occupied_cells": occupied_cells,
        "dominant_colors": dominant_colors,
        "element_count": counted,
    }


def _svg_element_bounds(element: ET.Element, *, name: str) -> tuple[float, float, float, float] | None:
    if name == "rect":
        x = _safe_float(element.attrib.get("x")) or 0.0
        y = _safe_float(element.attrib.get("y")) or 0.0
        width = _safe_float(element.attrib.get("width")) or 0.0
        height = _safe_float(element.attrib.get("height")) or 0.0
        return x, y, width, height
    if name == "circle":
        radius = _safe_float(element.attrib.get("r")) or 0.0
        cx = _safe_float(element.attrib.get("cx")) or 0.0
        cy = _safe_float(element.attrib.get("cy")) or 0.0
        return cx - radius, cy - radius, radius * 2, radius * 2
    if name == "ellipse":
        rx = _safe_float(element.attrib.get("rx")) or 0.0
        ry = _safe_float(element.attrib.get("ry")) or 0.0
        cx = _safe_float(element.attrib.get("cx")) or 0.0
        cy = _safe_float(element.attrib.get("cy")) or 0.0
        return cx - rx, cy - ry, rx * 2, ry * 2
    if name == "line":
        x1 = _safe_float(element.attrib.get("x1")) or 0.0
        y1 = _safe_float(element.attrib.get("y1")) or 0.0
        x2 = _safe_float(element.attrib.get("x2")) or 0.0
        y2 = _safe_float(element.attrib.get("y2")) or 0.0
        return min(x1, x2), min(y1, y2), max(abs(x2 - x1), 2.0), max(abs(y2 - y1), 2.0)
    if name in {"polyline", "polygon"}:
        return _points_bounds(str(element.attrib.get("points") or ""))
    if name == "text":
        text = " ".join(element.itertext()).strip()
        x = _safe_float(element.attrib.get("x")) or 0.0
        y = _safe_float(element.attrib.get("y")) or 0.0
        font_size = _element_font_size(element) or 16.0
        width = max(font_size * 1.2, min(font_size * 0.62 * max(len(text), 1), font_size * 80))
        height = max(font_size * 1.2, 8.0)
        return x, max(0.0, y - height), width, height
    return None


def _points_bounds(points: str) -> tuple[float, float, float, float] | None:
    numbers = [_safe_float(match.group(0)) for match in re.finditer(r"-?\d+(?:\.\d+)?", points)]
    clean = [number for number in numbers if number is not None]
    if len(clean) < 4:
        return None
    xs = clean[0::2]
    ys = clean[1::2]
    return min(xs), min(ys), max(xs) - min(xs), max(ys) - min(ys)


def _svg_element_color(element: ET.Element) -> str:
    for key in ("fill", "stroke"):
        value = str(element.attrib.get(key) or "").strip()
        if re.fullmatch(r"#[0-9A-Fa-f]{6}", value):
            return value.upper()
    style = str(element.attrib.get("style") or "")
    match = re.search(r"(?:fill|stroke)\s*:\s*(#[0-9A-Fa-f]{6})", style)
    if match:
        return match.group(1).upper()
    return "none"


def _element_font_size(element: ET.Element) -> float | None:
    direct = _safe_float(element.attrib.get("font-size"))
    if direct is not None:
        return direct
    style = str(element.attrib.get("style") or "")
    match = re.search(r"font-size\s*:\s*([0-9.]+)", style)
    return _safe_float(match.group(1)) if match else None


def _accumulate_render_grid(
    grid: list[float],
    *,
    x: float,
    y: float,
    width: float,
    height: float,
    view_width: float,
    view_height: float,
    weight: float,
) -> None:
    min_col = max(0, min(VISUAL_SIGNATURE_GRID_COLUMNS - 1, int((max(0.0, x) / view_width) * VISUAL_SIGNATURE_GRID_COLUMNS)))
    max_col = max(0, min(VISUAL_SIGNATURE_GRID_COLUMNS - 1, int((min(view_width, x + width) / view_width) * VISUAL_SIGNATURE_GRID_COLUMNS)))
    min_row = max(0, min(VISUAL_SIGNATURE_GRID_ROWS - 1, int((max(0.0, y) / view_height) * VISUAL_SIGNATURE_GRID_ROWS)))
    max_row = max(0, min(VISUAL_SIGNATURE_GRID_ROWS - 1, int((min(view_height, y + height) / view_height) * VISUAL_SIGNATURE_GRID_ROWS)))
    cell_count = max(1, (max_col - min_col + 1) * (max_row - min_row + 1))
    share = weight / cell_count
    for row in range(min_row, max_row + 1):
        for col in range(min_col, max_col + 1):
            grid[row * VISUAL_SIGNATURE_GRID_COLUMNS + col] += share


def _render_signature_delta(baseline: object, candidate: object) -> dict[str, object]:
    if not isinstance(baseline, dict) or not isinstance(candidate, dict):
        return {"available": False, "reason": "missing_signature"}
    baseline_grid = [float(value or 0) for value in baseline.get("grid") or []]
    candidate_grid = [float(value or 0) for value in candidate.get("grid") or []]
    if len(baseline_grid) != len(candidate_grid) or not baseline_grid:
        return {"available": False, "reason": "incompatible_grid"}
    changed_cells = sum(1 for left, right in zip(baseline_grid, candidate_grid, strict=False) if abs(left - right) > 0.015)
    moved_weight = sum(abs(left - right) for left, right in zip(baseline_grid, candidate_grid, strict=False)) / 2.0
    baseline_occupied = int(baseline.get("occupied_cells") or 0)
    candidate_occupied = int(candidate.get("occupied_cells") or 0)
    return {
        "available": bool(baseline.get("available")) and bool(candidate.get("available")),
        "grid_size": candidate.get("grid_size") or baseline.get("grid_size"),
        "changed_cells": changed_cells,
        "changed_cell_ratio": round(changed_cells / len(baseline_grid), 4),
        "moved_weight_ratio": round(moved_weight, 4),
        "baseline_occupied_cells": baseline_occupied,
        "candidate_occupied_cells": candidate_occupied,
        "occupied_cell_delta": candidate_occupied - baseline_occupied,
        "color_distribution_drift": _render_color_distribution_drift(baseline.get("dominant_colors"), candidate.get("dominant_colors")),
    }


def _render_color_distribution_drift(baseline: object, candidate: object) -> float:
    left = _render_color_weights(baseline)
    right = _render_color_weights(candidate)
    colors = set(left) | set(right)
    if not colors:
        return 0.0
    return round(sum(abs(left.get(color, 0.0) - right.get(color, 0.0)) for color in colors) / 2.0, 4)


def _render_color_weights(value: object) -> dict[str, float]:
    if not isinstance(value, list):
        return {}
    weights: dict[str, float] = {}
    for item in value:
        if not isinstance(item, dict):
            continue
        color = str(item.get("color") or "")
        if not color:
            continue
        weights[color] = float(item.get("weight") or 0)
    return weights


def _visual_diff_page_recommendations(
    *,
    deltas: dict[str, object],
    issues: list[PresentationVisualIssue],
) -> list[str]:
    issue_codes = {issue.code for issue in issues}
    recommendations: list[str] = []
    if "viewbox_changed" in issue_codes:
        recommendations.append("Regenerate the candidate page using the same canvas contract as the baseline.")
    if "text_density_regressed" in issue_codes:
        recommendations.append("Move added text to speaker notes or split dense content into another slide.")
    if "visual_complexity_regressed" in issue_codes:
        recommendations.append("Simplify decorative SVG elements before editable PPTX conversion.")
    if "render_layout_shift" in issue_codes or "render_occupancy_changed" in issue_codes:
        recommendations.append("Compare the rendered layout grid and restore the baseline page rhythm before accepting the candidate.")
    if "render_color_distribution_shift" in issue_codes:
        recommendations.append("Rebalance dominant fills so the candidate keeps the intended visual weight and contrast.")
    if "palette_regressed" in issue_codes or "new_palette_colors" in issue_codes:
        recommendations.append("Reduce new colors and re-anchor them to the locked theme palette.")
    if any(code.startswith("new_") for code in issue_codes):
        recommendations.append("Inspect the new candidate issue codes before treating the deck as improved.")
    if not recommendations and deltas.get("resolved_issue_codes"):
        recommendations.append("Candidate resolves prior issues without introducing detected regressions.")
    return recommendations[:6]


def _visual_diff_recommendations(
    pages: list[PresentationVisualDiffPage],
    *,
    error_count: int,
    warning_count: int,
) -> list[str]:
    if error_count:
        return ["Fix failed diff pages before delivery; the candidate removed pages or introduced hard regressions."]
    if warning_count:
        return ["Review warning diff pages for canvas, density, palette, or complexity regressions before accepting the candidate."]
    if any(page.deltas.get("resolved_issue_codes") for page in pages):
        return ["Candidate appears to improve prior visual issues without detected regressions."]
    return ["No material visual-review regressions detected by the current structural diff checks."]


def _annotations_from_visual_review(review: PresentationVisualReviewResult) -> list[PresentationAnnotation]:
    annotations: list[PresentationAnnotation] = []
    for page in review.pages:
        for issue_index, issue in enumerate(page.issues, start=1):
            annotations.append(
                PresentationAnnotation(
                    page_index=page.index,
                    page_name=page.name,
                    source="visual_review",
                    target="page",
                    severity=issue.severity,
                    code=issue.code,
                    message=issue.message,
                    region=_annotation_region_for_issue(issue.code, index=issue_index),
                    recommendation=_recommendation_for_issue(issue.code, page.recommendations),
                    path=page.path,
                )
            )
    return annotations


def _annotations_from_visual_diff(diff: PresentationVisualDiffResult) -> list[PresentationAnnotation]:
    annotations: list[PresentationAnnotation] = []
    for page in diff.pages:
        target_path = page.candidate_path or page.baseline_path
        target_name = page.candidate_name or page.baseline_name or f"page-{page.index}"
        target = "candidate" if page.candidate_path is not None else "baseline"
        for issue_index, issue in enumerate(page.issues, start=1):
            annotations.append(
                PresentationAnnotation(
                    page_index=page.index,
                    page_name=target_name,
                    source="visual_diff",
                    target=target,
                    severity=issue.severity,
                    code=issue.code,
                    message=issue.message,
                    region=_annotation_region_for_issue(issue.code, index=issue_index),
                    recommendation=_recommendation_for_issue(issue.code, page.recommendations),
                    path=target_path,
                )
            )
    return annotations


def _annotation_region_for_issue(code: str, *, index: int) -> dict[str, object]:
    normalized = code.removeprefix("new_")
    presets = {
        "viewbox_changed": {"x": 0.03, "y": 0.03, "width": 0.94, "height": 0.12, "anchor": "canvas"},
        "viewbox_contract_mismatch": {"x": 0.03, "y": 0.03, "width": 0.94, "height": 0.12, "anchor": "canvas"},
        "palette_contract_drift": {"x": 0.66, "y": 0.04, "width": 0.30, "height": 0.20, "anchor": "palette"},
        "palette_too_broad": {"x": 0.66, "y": 0.04, "width": 0.30, "height": 0.20, "anchor": "palette"},
        "palette_regressed": {"x": 0.66, "y": 0.04, "width": 0.30, "height": 0.20, "anchor": "palette"},
        "new_palette_colors": {"x": 0.66, "y": 0.04, "width": 0.30, "height": 0.20, "anchor": "palette"},
        "text_density_high": {"x": 0.08, "y": 0.22, "width": 0.84, "height": 0.52, "anchor": "text"},
        "text_density_regressed": {"x": 0.08, "y": 0.22, "width": 0.84, "height": 0.52, "anchor": "text"},
        "rhythm_density_mismatch": {"x": 0.08, "y": 0.22, "width": 0.84, "height": 0.52, "anchor": "rhythm"},
        "visual_complexity_high": {"x": 0.04, "y": 0.14, "width": 0.92, "height": 0.72, "anchor": "composition"},
        "visual_complexity_regressed": {"x": 0.04, "y": 0.14, "width": 0.92, "height": 0.72, "anchor": "composition"},
        "page_appears_blank": {"x": 0.20, "y": 0.20, "width": 0.60, "height": 0.60, "anchor": "page"},
        "page_added": {"x": 0.06, "y": 0.06, "width": 0.88, "height": 0.16, "anchor": "page"},
        "page_removed": {"x": 0.06, "y": 0.06, "width": 0.88, "height": 0.16, "anchor": "page"},
        "candidate_regressed_to_warning": {"x": 0.06, "y": 0.06, "width": 0.88, "height": 0.16, "anchor": "status"},
        "candidate_regressed_to_failed": {"x": 0.06, "y": 0.06, "width": 0.88, "height": 0.16, "anchor": "status"},
    }
    region = dict(presets.get(normalized, {"x": 0.08, "y": min(0.08 + (index - 1) * 0.10, 0.76), "width": 0.84, "height": 0.10, "anchor": "issue"}))
    region["coordinate_space"] = "normalized"
    return region


def _recommendation_for_issue(code: str, recommendations: tuple[str, ...] | list[str]) -> str | None:
    normalized = code.removeprefix("new_")
    for recommendation in recommendations:
        lower = recommendation.lower()
        if "viewbox" in normalized and "canvas" in lower:
            return recommendation
        if "palette" in normalized and ("color" in lower or "palette" in lower):
            return recommendation
        if "text" in normalized and ("text" in lower or "split" in lower):
            return recommendation
        if "complexity" in normalized and ("simplify" in lower or "svg" in lower):
            return recommendation
    return recommendations[0] if recommendations else None


def _annotation_pages_for_report(spec: PresentationAnnotationSpecResult) -> list[dict[str, object]]:
    pages: dict[tuple[int, str], dict[str, object]] = {}
    for annotation in spec.annotations:
        key = (annotation.page_index, str(annotation.path or annotation.page_name))
        pages.setdefault(
            key,
            {
                "index": annotation.page_index,
                "name": annotation.page_name,
                "path": annotation.path,
                "target": annotation.target,
                "annotations": [],
            },
        )["annotations"].append(annotation)
    return [pages[key] for key in sorted(pages)]


def _annotation_manifest_payload(annotation: PresentationAnnotation, *, spec: PresentationAnnotationSpecResult) -> dict[str, object]:
    payload = annotation.to_payload()
    if annotation.path is not None:
        root = spec.candidate_path or spec.target_path or spec.baseline_path or annotation.path.parent
        payload["path"] = _artifact_display_path(annotation.path, root=root)
    return payload


def _annotation_page_manifest_payload(page: dict[str, object]) -> dict[str, object]:
    annotations = page.get("annotations") if isinstance(page.get("annotations"), list) else []
    raw_path = page.get("path")
    return {
        "index": page.get("index"),
        "name": page.get("name"),
        "target": page.get("target"),
        "path": raw_path.name if isinstance(raw_path, Path) else None,
        "annotation_count": len(annotations),
    }


def _xml_local_name(tag: object) -> str:
    text = str(tag)
    if "}" in text:
        return text.rsplit("}", 1)[1]
    return text


def _viewbox_size(viewbox: object) -> tuple[float | None, float | None]:
    if not isinstance(viewbox, str):
        return None, None
    parts = [part for part in re.split(r"[\s,]+", viewbox.strip()) if part]
    if len(parts) != 4:
        return None, None
    return _safe_float(parts[2]), _safe_float(parts[3])


def _safe_float(value: object) -> float | None:
    try:
        return float(value)
    except (TypeError, ValueError):
        return None


def _body_to_bullets(lines: list[str]) -> list[str]:
    bullets: list[str] = []
    paragraph_buffer: list[str] = []
    for raw_line in lines:
        line = raw_line.strip()
        if not line:
            _flush_paragraph(paragraph_buffer, bullets)
            continue
        bullet = re.match(r"^(?:[-*+]|\d+[.)])\s+(.+)$", line)
        if bullet:
            _flush_paragraph(paragraph_buffer, bullets)
            bullets.append(bullet.group(1).strip()[:220])
            continue
        paragraph_buffer.append(line)
    _flush_paragraph(paragraph_buffer, bullets)
    return [bullet for bullet in bullets if bullet]


def _flush_paragraph(buffer: list[str], bullets: list[str]) -> None:
    if not buffer:
        return
    paragraph = " ".join(buffer).strip()
    buffer.clear()
    if paragraph:
        bullets.append(paragraph[:220])


def _plain_chunks(content: str) -> list[str]:
    chunks = [line.strip() for line in content.splitlines() if line.strip()]
    return chunks[:MAX_BULLETS_PER_SLIDE] or ["No content provided."]


def _read_template_metadata(design_spec_path: Path | None) -> dict[str, object]:
    if design_spec_path is None:
        return {}
    text = design_spec_path.read_text(encoding="utf-8")
    frontmatter, body = _split_simple_frontmatter(text)
    metadata: dict[str, object] = dict(frontmatter)
    if "summary" not in metadata:
        overview = _extract_markdown_section(body, "Template Overview")
        if overview:
            metadata["summary"] = _first_nonempty_line(overview)
    if "primary_color" not in metadata:
        color = _first_hex_color(text)
        if color:
            metadata["primary_color"] = color
    if "replication_mode" not in metadata:
        mode_match = re.search(r"replication[_ -]mode\s*[:：]\s*([A-Za-z0-9_-]+)", text, re.IGNORECASE)
        if mode_match:
            metadata["replication_mode"] = mode_match.group(1)
    return metadata


def _split_simple_frontmatter(text: str) -> tuple[dict[str, object], str]:
    if not text.startswith("---\n"):
        return {}, text
    end = text.find("\n---\n", 4)
    if end == -1:
        return {}, text
    block = text[4:end]
    body = text[end + 5:]
    data: dict[str, object] = {}
    for raw_line in block.splitlines():
        line = raw_line.strip()
        if not line or line.startswith("#") or ":" not in line:
            continue
        key, raw_value = line.split(":", 1)
        key = key.strip()
        value = raw_value.strip().strip('"').strip("'")
        if value.startswith("[") and value.endswith("]"):
            items = [item.strip().strip('"').strip("'") for item in value[1:-1].split(",")]
            data[key] = [item for item in items if item]
        elif "," in value and key in {"keywords", "use_cases"}:
            data[key] = [item.strip() for item in value.split(",") if item.strip()]
        else:
            data[key] = value
    return data, body


def _extract_markdown_section(text: str, title: str) -> str:
    pattern = re.compile(rf"^##\s+(?:[IVX]+\.\s+)?{re.escape(title)}\b(.*?)(?=^##\s+|\Z)", re.MULTILINE | re.DOTALL)
    match = pattern.search(text)
    return match.group(1).strip() if match else ""


def _first_nonempty_line(text: str) -> str | None:
    for line in text.splitlines():
        stripped = re.sub(r"^[|#*\-\s]+", "", line).strip()
        stripped = re.sub(r"[|*`\s]+$", "", stripped).strip()
        if stripped and not stripped.startswith("---"):
            return stripped[:240]
    return None


def _first_hex_color(text: str) -> str | None:
    match = re.search(r"#[0-9A-Fa-f]{6}", text)
    return match.group(0).upper() if match else None


def _safe_path_segment(value: str) -> str:
    cleaned = re.sub(r"[^A-Za-z0-9._-]+", "-", value.strip()).strip(".-")
    return cleaned[:80] or "template"


def _artifact_display_path(path: Path, *, root: Path) -> str:
    try:
        base = root if root.is_dir() else root.parent
        return path.relative_to(base).as_posix()
    except ValueError:
        return path.name


def _visual_diff_page_manifest(
    page: PresentationVisualDiffPage,
    *,
    baseline_root: Path,
    candidate_root: Path,
) -> dict[str, object]:
    return {
        "index": page.index,
        "status": page.status,
        "baseline_name": page.baseline_name,
        "candidate_name": page.candidate_name,
        "baseline_path": _artifact_display_path(page.baseline_path, root=baseline_root) if page.baseline_path else None,
        "candidate_path": _artifact_display_path(page.candidate_path, root=candidate_root) if page.candidate_path else None,
        "baseline_status": page.baseline_status,
        "candidate_status": page.candidate_status,
        "deltas": dict(page.deltas),
        "issues": [issue.to_payload() for issue in page.issues],
        "recommendations": list(page.recommendations),
    }


def _safe_svg_preview_data_uri(path: Path | None, *, label: str) -> tuple[str | None, list[str]]:
    if path is None:
        return None, []
    if not path.exists():
        return None, [f"{label} SVG is missing"]
    size_bytes = path.stat().st_size
    if size_bytes > MAX_TEMPLATE_PREVIEW_SVG_BYTES:
        return None, [f"{label} SVG is too large for inline preview"]
    return _svg_data_uri(path), []


def _svg_data_uri(path: Path) -> str:
    raw = path.read_bytes()
    encoded = base64.b64encode(raw).decode("ascii")
    return f"data:image/svg+xml;base64,{encoded}"


def _compact_preflight_payload(payload: dict[str, object]) -> dict[str, object]:
    issues = payload.get("issues")
    return {
        "status": payload.get("status"),
        "error_count": payload.get("error_count", 0),
        "warning_count": payload.get("warning_count", 0),
        "issues": list(issues or [])[:12] if isinstance(issues, list) else [],
        "summary": dict(payload.get("summary") or {}),
    }


def _template_preview_status(pages: list[dict[str, object]], warnings: list[str]) -> str:
    failed = False
    for page in pages:
        if page.get("warnings"):
            failed = True
        preflight = page.get("preflight")
        if isinstance(preflight, dict) and preflight.get("status") == "failed":
            failed = True
    if failed:
        return "failed"
    if warnings:
        return "warning"
    return "passed"


def _render_template_preview_html(
    *,
    template_id: str,
    metadata: dict[str, object],
    pages: list[dict[str, object]],
    status: str,
    warnings: list[str],
) -> str:
    escaped_title = html.escape(template_id)
    summary = html.escape(str(metadata.get("summary") or "No summary provided."))
    primary = html.escape(str(metadata.get("primary_color") or "#3A506B"))
    status_class = "failed" if status == "failed" else "warning" if status == "warning" else "passed"
    warnings_html = "".join(f"<li>{html.escape(warning)}</li>" for warning in warnings) or "<li>No template-level warnings.</li>"
    pages_html = "\n".join(_render_template_preview_page(page) for page in pages)
    metadata_html = html.escape(json.dumps(metadata, ensure_ascii=False, indent=2))
    return f"""<!doctype html>
<html lang="en">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>{escaped_title} preview</title>
<style>
:root {{
  color-scheme: light;
  --ink: #203247;
  --muted: #66788A;
  --line: #D8E3EB;
  --panel: #FFFFFF;
  --wash: #F5FAFC;
  --brand: {primary};
  --accent: #5BC0BE;
}}
* {{ box-sizing: border-box; }}
body {{
  margin: 0;
  padding: 28px;
  background: var(--wash);
  color: var(--ink);
  font-family: Inter, -apple-system, BlinkMacSystemFont, "Segoe UI", sans-serif;
}}
header {{
  max-width: 1180px;
  margin: 0 auto 20px;
}}
h1 {{
  margin: 0 0 8px;
  font-size: 24px;
  letter-spacing: 0;
}}
.summary {{
  margin: 0;
  color: var(--muted);
  line-height: 1.5;
}}
.badge {{
  display: inline-flex;
  align-items: center;
  margin: 14px 0;
  padding: 4px 10px;
  border: 1px solid var(--line);
  border-radius: 999px;
  font-size: 12px;
  font-weight: 600;
  text-transform: uppercase;
}}
.badge.passed {{ color: #107A52; background: #E7F7EF; }}
.badge.warning {{ color: #8A5A00; background: #FFF5D6; }}
.badge.failed {{ color: #B42318; background: #FDECEC; }}
.grid {{
  max-width: 1180px;
  margin: 0 auto;
  display: grid;
  grid-template-columns: repeat(auto-fit, minmax(320px, 1fr));
  gap: 18px;
}}
.page {{
  background: var(--panel);
  border: 1px solid var(--line);
  border-radius: 10px;
  overflow: hidden;
  box-shadow: 0 14px 30px rgba(32, 50, 71, 0.07);
}}
.frame {{
  background: #EDF5F8;
  border-bottom: 1px solid var(--line);
  aspect-ratio: 16 / 9;
  display: grid;
  place-items: center;
  overflow: hidden;
}}
.frame img {{
  width: 100%;
  height: 100%;
  object-fit: contain;
  display: block;
}}
.missing {{
  color: var(--muted);
  font-size: 13px;
}}
.meta {{
  padding: 12px 14px 14px;
}}
.meta h2 {{
  margin: 0 0 6px;
  font-size: 14px;
}}
.meta p {{
  margin: 3px 0;
  color: var(--muted);
  font-size: 12px;
  line-height: 1.45;
}}
.issues, .warnings {{
  max-width: 1180px;
  margin: 18px auto 0;
  background: var(--panel);
  border: 1px solid var(--line);
  border-radius: 10px;
  padding: 14px 18px;
}}
pre {{
  white-space: pre-wrap;
  overflow-wrap: anywhere;
  color: var(--muted);
  font-size: 12px;
}}
</style>
</head>
<body>
<header>
  <h1>{escaped_title}</h1>
  <p class="summary">{summary}</p>
  <div class="badge {status_class}">{html.escape(status)}</div>
</header>
<main class="grid">
{pages_html}
</main>
<section class="warnings">
  <h2>Warnings</h2>
  <ul>{warnings_html}</ul>
</section>
<section class="issues">
  <h2>Metadata</h2>
  <pre>{metadata_html}</pre>
</section>
</body>
</html>
"""


def _render_template_preview_page(page: dict[str, object]) -> str:
    name = html.escape(str(page.get("name") or "Untitled"))
    data_uri = page.get("preview_data_uri")
    image_html = (
        f'<img alt="{name}" src="{html.escape(str(data_uri), quote=True)}">'
        if isinstance(data_uri, str) and data_uri
        else '<div class="missing">Preview omitted for this SVG.</div>'
    )
    warnings = page.get("warnings")
    warning_text = ", ".join(str(item) for item in warnings or []) or "none"
    preflight = page.get("preflight") if isinstance(page.get("preflight"), dict) else {}
    preflight_status = str(preflight.get("status") or "not checked")
    issues = preflight.get("issues")
    issue_text = "; ".join(str(issue.get("code") or issue.get("message")) for issue in issues or [] if isinstance(issue, dict)) or "none"
    return f"""<article class="page">
  <div class="frame">{image_html}</div>
  <div class="meta">
    <h2>{name}</h2>
    <p>viewBox: {html.escape(str(page.get("viewBox") or "missing"))}</p>
    <p>size: {html.escape(str(page.get("size_bytes") or 0))} bytes</p>
    <p>preflight: {html.escape(preflight_status)}</p>
    <p>warnings: {html.escape(warning_text)}</p>
    <p>issues: {html.escape(issue_text)}</p>
  </div>
</article>"""


def _render_visual_review_report_html(
    *,
    title: str,
    status: str,
    summary: dict[str, object],
    recommendations: list[str],
    pages: list[dict[str, object]],
    warnings: list[str],
) -> str:
    escaped_title = html.escape(title)
    status_class = "failed" if status == "failed" else "warning" if status == "warning" else "passed"
    summary_items = [
        ("errors", summary.get("error_count", 0)),
        ("warnings", summary.get("warning_count", 0)),
        ("text chars", summary.get("total_text_characters", 0)),
        ("text nodes", summary.get("total_text_nodes", 0)),
        ("contract slides", summary.get("contract_slide_count", 0)),
    ]
    summary_html = "".join(
        f"<div><strong>{html.escape(str(value))}</strong><span>{html.escape(label)}</span></div>"
        for label, value in summary_items
    )
    palette = summary.get("palette") if isinstance(summary.get("palette"), list) else []
    palette_html = "".join(f'<i style="background:{html.escape(str(color))}"></i>' for color in palette) or "<span>No palette detected.</span>"
    recommendations_html = "".join(f"<li>{html.escape(item)}</li>" for item in recommendations) or "<li>No deck-level recommendations.</li>"
    warnings_html = "".join(f"<li>{html.escape(item)}</li>" for item in warnings) or "<li>No report-level warnings.</li>"
    pages_html = "\n".join(_render_visual_review_report_page(page) for page in pages)
    return f"""<!doctype html>
<html lang="en">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>{escaped_title}</title>
<style>
:root {{
  color-scheme: light;
  --ink: #1F2A37;
  --muted: #66788A;
  --line: #D8E3EB;
  --panel: #FFFFFF;
  --wash: #F5FAFC;
  --accent: #5BC0BE;
  --brand: #3A506B;
  --danger: #B42318;
  --warn: #8A5A00;
  --ok: #107A52;
}}
* {{ box-sizing: border-box; }}
body {{
  margin: 0;
  padding: 28px;
  background: var(--wash);
  color: var(--ink);
  font-family: Inter, -apple-system, BlinkMacSystemFont, "Segoe UI", sans-serif;
}}
header, main, section {{
  max-width: 1220px;
  margin-left: auto;
  margin-right: auto;
}}
header {{
  margin-bottom: 18px;
}}
h1 {{
  margin: 0 0 8px;
  font-size: 24px;
  letter-spacing: 0;
}}
h2 {{
  margin: 0 0 10px;
  font-size: 15px;
  letter-spacing: 0;
}}
.badge {{
  display: inline-flex;
  align-items: center;
  margin: 8px 0 12px;
  padding: 4px 10px;
  border-radius: 999px;
  font-size: 12px;
  font-weight: 650;
  text-transform: uppercase;
}}
.badge.passed {{ color: var(--ok); background: #E7F7EF; }}
.badge.warning {{ color: var(--warn); background: #FFF5D6; }}
.badge.failed {{ color: var(--danger); background: #FDECEC; }}
.summary-grid {{
  display: grid;
  grid-template-columns: repeat(auto-fit, minmax(120px, 1fr));
  gap: 10px;
  margin: 12px 0 16px;
}}
.summary-grid div, .panel {{
  background: var(--panel);
  border: 1px solid var(--line);
  border-radius: 10px;
  box-shadow: 0 12px 30px rgba(32, 50, 71, 0.06);
}}
.summary-grid div {{
  padding: 12px 14px;
}}
.summary-grid strong {{
  display: block;
  font-size: 20px;
  line-height: 1.1;
}}
.summary-grid span, .muted {{
  color: var(--muted);
  font-size: 12px;
}}
.palette {{
  display: flex;
  align-items: center;
  gap: 6px;
  margin: 8px 0 14px;
  min-height: 18px;
}}
.palette i {{
  width: 18px;
  height: 18px;
  border-radius: 50%;
  border: 1px solid rgba(0,0,0,0.12);
  display: inline-block;
}}
.panel {{
  padding: 14px 18px;
  margin-bottom: 16px;
}}
.deck {{
  display: grid;
  grid-template-columns: 1fr;
  gap: 16px;
}}
.page {{
  display: grid;
  grid-template-columns: minmax(260px, 42%) 1fr;
  overflow: hidden;
}}
.frame {{
  background: #EDF5F8;
  border-right: 1px solid var(--line);
  min-height: 220px;
  display: grid;
  place-items: center;
}}
.frame img {{
  width: 100%;
  height: 100%;
  object-fit: contain;
  display: block;
}}
.missing {{
  color: var(--muted);
  font-size: 13px;
}}
.details {{
  padding: 14px 16px 16px;
}}
.details header {{
  margin: 0 0 10px;
  display: flex;
  align-items: center;
  justify-content: space-between;
  gap: 10px;
}}
.details h3 {{
  margin: 0;
  font-size: 15px;
}}
.details code {{
  background: #EEF5F7;
  border-radius: 6px;
  padding: 2px 6px;
  font-size: 12px;
}}
.issue {{
  border-top: 1px solid var(--line);
  padding: 9px 0;
}}
.issue strong {{
  display: block;
  font-size: 13px;
}}
.issue p, .rec li, .metrics li {{
  margin: 4px 0;
  color: var(--muted);
  font-size: 12px;
  line-height: 1.45;
}}
.rec, .metrics {{
  padding-left: 18px;
  margin: 6px 0 0;
}}
@media (max-width: 760px) {{
  body {{ padding: 18px; }}
  .page {{ grid-template-columns: 1fr; }}
  .frame {{ border-right: 0; border-bottom: 1px solid var(--line); }}
}}
</style>
</head>
<body>
<header>
  <h1>{escaped_title}</h1>
  <div class="badge {status_class}">{html.escape(status)}</div>
  <div class="summary-grid">{summary_html}</div>
  <div class="palette">{palette_html}</div>
</header>
<section class="panel">
  <h2>Recommendations</h2>
  <ul>{recommendations_html}</ul>
</section>
<main class="deck">
{pages_html}
</main>
<section class="panel">
  <h2>Report Warnings</h2>
  <ul>{warnings_html}</ul>
</section>
</body>
</html>
"""


def _render_visual_review_report_page(page: dict[str, object]) -> str:
    name = html.escape(str(page.get("name") or "Untitled"))
    status = str(page.get("status") or "unknown")
    status_class = "failed" if status == "failed" else "warning" if status == "warning" else "passed"
    data_uri = page.get("preview_data_uri")
    image_html = (
        f'<img alt="{name}" src="{html.escape(str(data_uri), quote=True)}">'
        if isinstance(data_uri, str) and data_uri
        else '<div class="missing">Preview omitted for this SVG.</div>'
    )
    issues = page.get("issues") if isinstance(page.get("issues"), list) else []
    issue_html = "".join(_render_visual_review_issue(issue) for issue in issues if isinstance(issue, dict)) or '<div class="issue"><strong>No issues</strong><p>Page passed the current checks.</p></div>'
    recommendations = page.get("recommendations") if isinstance(page.get("recommendations"), list) else []
    recommendations_html = "".join(f"<li>{html.escape(str(item))}</li>" for item in recommendations) or "<li>No page-level recommendations.</li>"
    metrics = page.get("metrics") if isinstance(page.get("metrics"), dict) else {}
    metrics_html = "".join(
        f"<li>{html.escape(label)}: {html.escape(str(value))}</li>"
        for label, value in (
            ("text characters", metrics.get("text_characters", 0)),
            ("text nodes", metrics.get("text_node_count", 0)),
            ("shapes", metrics.get("shape_count", 0)),
            ("colors", metrics.get("color_count", 0)),
            ("elements", metrics.get("element_count", 0)),
        )
    )
    return f"""<article class="panel page">
  <div class="frame">{image_html}</div>
  <div class="details">
    <header>
      <h3>{name}</h3>
      <span class="badge {status_class}">{html.escape(status)}</span>
    </header>
    <p class="muted">path: <code>{html.escape(str(page.get("path") or ""))}</code></p>
    <p class="muted">viewBox: <code>{html.escape(str(page.get("viewBox") or "missing"))}</code></p>
    <ul class="metrics">{metrics_html}</ul>
    {issue_html}
    <h2>Next edits</h2>
    <ul class="rec">{recommendations_html}</ul>
  </div>
</article>"""


def _render_visual_review_issue(issue: dict[str, object]) -> str:
    severity = html.escape(str(issue.get("severity") or "warning"))
    code = html.escape(str(issue.get("code") or "issue"))
    message = html.escape(str(issue.get("message") or ""))
    return f"""<div class="issue">
  <strong>{severity}: {code}</strong>
  <p>{message}</p>
</div>"""


def _render_visual_diff_report_html(
    *,
    title: str,
    status: str,
    summary: dict[str, object],
    recommendations: list[str],
    pages: list[dict[str, object]],
    warnings: list[str],
) -> str:
    escaped_title = html.escape(title)
    status_class = "failed" if status == "failed" else "warning" if status == "warning" else "passed"
    summary_items = [
        ("changed pages", summary.get("changed_pages", 0)),
        ("errors", summary.get("error_count", 0)),
        ("warnings", summary.get("warning_count", 0)),
        ("baseline pages", summary.get("baseline_page_count", 0)),
        ("candidate pages", summary.get("candidate_page_count", 0)),
    ]
    summary_html = "".join(
        f"<div><strong>{html.escape(str(value))}</strong><span>{html.escape(label)}</span></div>"
        for label, value in summary_items
    )
    recommendations_html = "".join(f"<li>{html.escape(item)}</li>" for item in recommendations) or "<li>No deck-level recommendations.</li>"
    warnings_html = "".join(f"<li>{html.escape(item)}</li>" for item in warnings) or "<li>No report-level warnings.</li>"
    pages_html = "\n".join(_render_visual_diff_report_page(page) for page in pages)
    return f"""<!doctype html>
<html lang="en">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>{escaped_title}</title>
<style>
:root {{
  color-scheme: light;
  --ink: #1F2A37;
  --muted: #66788A;
  --line: #D8E3EB;
  --panel: #FFFFFF;
  --wash: #F5FAFC;
  --danger: #B42318;
  --warn: #8A5A00;
  --ok: #107A52;
}}
* {{ box-sizing: border-box; }}
body {{
  margin: 0;
  padding: 28px;
  background: var(--wash);
  color: var(--ink);
  font-family: Inter, -apple-system, BlinkMacSystemFont, "Segoe UI", sans-serif;
}}
header, main, section {{
  max-width: 1280px;
  margin-left: auto;
  margin-right: auto;
}}
h1 {{
  margin: 0 0 8px;
  font-size: 24px;
  letter-spacing: 0;
}}
h2 {{
  margin: 0 0 10px;
  font-size: 15px;
}}
.badge {{
  display: inline-flex;
  align-items: center;
  padding: 4px 10px;
  border-radius: 999px;
  font-size: 12px;
  font-weight: 650;
  text-transform: uppercase;
}}
.badge.passed, .badge.unchanged {{ color: var(--ok); background: #E7F7EF; }}
.badge.warning, .badge.added {{ color: var(--warn); background: #FFF5D6; }}
.badge.failed, .badge.removed {{ color: var(--danger); background: #FDECEC; }}
.summary-grid {{
  display: grid;
  grid-template-columns: repeat(auto-fit, minmax(130px, 1fr));
  gap: 10px;
  margin: 14px 0 16px;
}}
.summary-grid div, .panel {{
  background: var(--panel);
  border: 1px solid var(--line);
  border-radius: 10px;
  box-shadow: 0 12px 30px rgba(32, 50, 71, 0.06);
}}
.summary-grid div {{
  padding: 12px 14px;
}}
.summary-grid strong {{
  display: block;
  font-size: 20px;
  line-height: 1.1;
}}
.summary-grid span, .muted {{
  color: var(--muted);
  font-size: 12px;
}}
.panel {{
  padding: 14px 18px;
  margin-bottom: 16px;
}}
.deck {{
  display: grid;
  grid-template-columns: 1fr;
  gap: 16px;
}}
.page {{
  overflow: hidden;
}}
.page-head {{
  display: flex;
  align-items: center;
  justify-content: space-between;
  gap: 10px;
  margin-bottom: 10px;
}}
.frames {{
  display: grid;
  grid-template-columns: 1fr 1fr;
  border: 1px solid var(--line);
  border-radius: 10px;
  overflow: hidden;
}}
.frame {{
  background: #EDF5F8;
  min-height: 220px;
  display: grid;
  place-items: center;
}}
.frame:first-child {{
  border-right: 1px solid var(--line);
}}
.frame img {{
  width: 100%;
  height: 100%;
  object-fit: contain;
  display: block;
}}
.frame b {{
  position: absolute;
  align-self: start;
  justify-self: start;
  margin: 8px;
  background: rgba(255,255,255,0.82);
  border: 1px solid var(--line);
  border-radius: 999px;
  padding: 2px 8px;
  font-size: 11px;
}}
.missing {{
  color: var(--muted);
  font-size: 13px;
}}
.issue {{
  border-top: 1px solid var(--line);
  padding: 9px 0;
}}
.issue strong, .delta strong {{
  display: block;
  font-size: 13px;
}}
.issue p, .delta li, .rec li {{
  margin: 4px 0;
  color: var(--muted);
  font-size: 12px;
  line-height: 1.45;
}}
.delta, .rec {{
  padding-left: 18px;
  margin: 6px 0 0;
}}
@media (max-width: 760px) {{
  body {{ padding: 18px; }}
  .frames {{ grid-template-columns: 1fr; }}
  .frame:first-child {{ border-right: 0; border-bottom: 1px solid var(--line); }}
}}
</style>
</head>
<body>
<header>
  <h1>{escaped_title}</h1>
  <div class="badge {status_class}">{html.escape(status)}</div>
  <div class="summary-grid">{summary_html}</div>
</header>
<section class="panel">
  <h2>Recommendations</h2>
  <ul>{recommendations_html}</ul>
</section>
<main class="deck">
{pages_html}
</main>
<section class="panel">
  <h2>Report Warnings</h2>
  <ul>{warnings_html}</ul>
</section>
</body>
</html>
"""


def _render_visual_diff_report_page(page: dict[str, object]) -> str:
    status = str(page.get("status") or "unknown")
    status_class = "failed" if status == "failed" else "warning" if status == "warning" else status
    title = html.escape(str(page.get("candidate_name") or page.get("baseline_name") or f"Page {page.get('index')}"))
    baseline_preview = page.get("baseline_preview_data_uri")
    candidate_preview = page.get("candidate_preview_data_uri")
    baseline_html = _render_visual_diff_frame(label="Baseline", name=page.get("baseline_path"), data_uri=baseline_preview)
    candidate_html = _render_visual_diff_frame(label="Candidate", name=page.get("candidate_path"), data_uri=candidate_preview)
    issues = page.get("issues") if isinstance(page.get("issues"), list) else []
    issue_html = "".join(_render_visual_review_issue(issue) for issue in issues if isinstance(issue, dict)) or '<div class="issue"><strong>No regressions</strong><p>No page-level regression detected.</p></div>'
    deltas = page.get("deltas") if isinstance(page.get("deltas"), dict) else {}
    delta_html = _render_visual_diff_delta_list(deltas)
    recommendations = page.get("recommendations") if isinstance(page.get("recommendations"), list) else []
    rec_html = "".join(f"<li>{html.escape(str(item))}</li>" for item in recommendations) or "<li>No page-level recommendations.</li>"
    return f"""<article class="panel page">
  <div class="page-head">
    <h2>{title}</h2>
    <span class="badge {html.escape(status_class)}">{html.escape(status)}</span>
  </div>
  <div class="frames">{baseline_html}{candidate_html}</div>
  <ul class="delta">{delta_html}</ul>
  {issue_html}
  <h2>Next edits</h2>
  <ul class="rec">{rec_html}</ul>
</article>"""


def _render_visual_diff_frame(*, label: str, name: object, data_uri: object) -> str:
    caption = html.escape(str(name or "missing"))
    image_html = (
        f'<img alt="{html.escape(label)}" src="{html.escape(str(data_uri), quote=True)}">'
        if isinstance(data_uri, str) and data_uri
        else '<div class="missing">Preview omitted.</div>'
    )
    return f"""<div class="frame">
  <b>{html.escape(label)}: {caption}</b>
  {image_html}
</div>"""


def _render_visual_diff_delta_list(deltas: dict[str, object]) -> str:
    fields = ["text_characters", "text_node_count", "shape_count", "element_count", "color_count"]
    items: list[str] = []
    for field in fields:
        payload = deltas.get(field)
        if isinstance(payload, dict):
            items.append(
                f"<li>{html.escape(field)}: {html.escape(str(payload.get('baseline', 0)))} -> "
                f"{html.escape(str(payload.get('candidate', 0)))} "
                f"({html.escape(str(payload.get('delta', 0)))})</li>"
            )
    for field in ("new_issue_codes", "resolved_issue_codes", "colors_added"):
        value = deltas.get(field)
        if isinstance(value, list) and value:
            items.append(f"<li>{html.escape(field)}: {html.escape(', '.join(str(item) for item in value[:10]))}</li>")
    render_delta = deltas.get("render_signature")
    if isinstance(render_delta, dict) and render_delta.get("available"):
        items.append(
            "<li>render_signature: "
            f"{html.escape(str(render_delta.get('changed_cells', 0)))} cells changed, "
            f"{html.escape(str(render_delta.get('moved_weight_ratio', 0)))} moved weight, "
            f"{html.escape(str(render_delta.get('color_distribution_drift', 0)))} color drift</li>"
        )
    return "".join(items) or "<li>No material metric deltas.</li>"


def _render_annotation_report_html(
    *,
    title: str,
    status: str,
    mode: str,
    summary: dict[str, object],
    recommendations: list[str],
    pages: list[dict[str, object]],
    annotations: list[PresentationAnnotation],
    warnings: list[str],
) -> str:
    escaped_title = html.escape(title)
    status_class = "failed" if status == "failed" else "warning" if status == "warning" else "passed"
    summary_items = [
        ("annotations", len(annotations)),
        ("errors", summary.get("error_count", 0)),
        ("warnings", summary.get("warning_count", 0)),
        ("mode", mode),
    ]
    summary_html = "".join(
        f"<div><strong>{html.escape(str(value))}</strong><span>{html.escape(label)}</span></div>"
        for label, value in summary_items
    )
    recommendations_html = "".join(f"<li>{html.escape(item)}</li>" for item in recommendations) or "<li>No deck-level recommendations.</li>"
    warnings_html = "".join(f"<li>{html.escape(item)}</li>" for item in warnings) or "<li>No report-level warnings.</li>"
    pages_html = "\n".join(_render_annotation_report_page(page) for page in pages)
    return f"""<!doctype html>
<html lang="en">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>{escaped_title}</title>
<style>
:root {{
  color-scheme: light;
  --ink: #1F2A37;
  --muted: #66788A;
  --line: #D8E3EB;
  --panel: #FFFFFF;
  --wash: #F5FAFC;
  --danger: #B42318;
  --warn: #8A5A00;
  --ok: #107A52;
}}
* {{ box-sizing: border-box; }}
body {{
  margin: 0;
  padding: 28px;
  background: var(--wash);
  color: var(--ink);
  font-family: Inter, -apple-system, BlinkMacSystemFont, "Segoe UI", sans-serif;
}}
header, main, section {{
  max-width: 1220px;
  margin-left: auto;
  margin-right: auto;
}}
h1 {{ margin: 0 0 8px; font-size: 24px; letter-spacing: 0; }}
h2 {{ margin: 0 0 10px; font-size: 15px; }}
.badge {{
  display: inline-flex;
  align-items: center;
  padding: 4px 10px;
  border-radius: 999px;
  font-size: 12px;
  font-weight: 650;
  text-transform: uppercase;
}}
.badge.passed {{ color: var(--ok); background: #E7F7EF; }}
.badge.warning {{ color: var(--warn); background: #FFF5D6; }}
.badge.failed {{ color: var(--danger); background: #FDECEC; }}
.summary-grid {{
  display: grid;
  grid-template-columns: repeat(auto-fit, minmax(130px, 1fr));
  gap: 10px;
  margin: 14px 0 16px;
}}
.summary-grid div, .panel {{
  background: var(--panel);
  border: 1px solid var(--line);
  border-radius: 10px;
  box-shadow: 0 12px 30px rgba(32, 50, 71, 0.06);
}}
.summary-grid div {{ padding: 12px 14px; }}
.summary-grid strong {{ display: block; font-size: 20px; line-height: 1.1; }}
.summary-grid span, .muted {{ color: var(--muted); font-size: 12px; }}
.panel {{ padding: 14px 18px; margin-bottom: 16px; }}
.deck {{ display: grid; grid-template-columns: 1fr; gap: 16px; }}
.page-grid {{ display: grid; grid-template-columns: minmax(280px, 52%) 1fr; gap: 14px; }}
.stage {{
  position: relative;
  background: #EDF5F8;
  border: 1px solid var(--line);
  border-radius: 10px;
  aspect-ratio: 16 / 9;
  overflow: hidden;
}}
.stage img {{ width: 100%; height: 100%; object-fit: contain; display: block; }}
.missing {{ display: grid; place-items: center; height: 100%; color: var(--muted); font-size: 13px; }}
.mark {{
  position: absolute;
  border: 2px solid var(--warn);
  background: rgba(91, 192, 190, 0.16);
  border-radius: 8px;
  min-width: 24px;
  min-height: 20px;
}}
.mark.error {{ border-color: var(--danger); background: rgba(180, 35, 24, 0.14); }}
.mark.warning {{ border-color: var(--warn); }}
.mark-label {{
  position: absolute;
  top: -10px;
  left: -2px;
  background: #fff;
  border: 1px solid currentColor;
  border-radius: 999px;
  padding: 1px 6px;
  font-size: 10px;
  font-weight: 700;
}}
.annotation {{
  border-top: 1px solid var(--line);
  padding: 9px 0;
}}
.annotation strong {{ display: block; font-size: 13px; }}
.annotation p, .rec li {{
  margin: 4px 0;
  color: var(--muted);
  font-size: 12px;
  line-height: 1.45;
}}
.rec {{ padding-left: 18px; margin: 6px 0 0; }}
@media (max-width: 820px) {{
  body {{ padding: 18px; }}
  .page-grid {{ grid-template-columns: 1fr; }}
}}
</style>
</head>
<body>
<header>
  <h1>{escaped_title}</h1>
  <div class="badge {status_class}">{html.escape(status)}</div>
  <div class="summary-grid">{summary_html}</div>
</header>
<section class="panel">
  <h2>Recommendations</h2>
  <ul>{recommendations_html}</ul>
</section>
<main class="deck">
{pages_html}
</main>
<section class="panel">
  <h2>Report Warnings</h2>
  <ul>{warnings_html}</ul>
</section>
</body>
</html>
"""


def _render_annotation_report_page(page: dict[str, object]) -> str:
    name = html.escape(str(page.get("name") or "Untitled"))
    preview_data_uri = page.get("preview_data_uri")
    annotations = page.get("annotations") if isinstance(page.get("annotations"), list) else []
    image_html = (
        f'<img alt="{name}" src="{html.escape(str(preview_data_uri), quote=True)}">'
        if isinstance(preview_data_uri, str) and preview_data_uri
        else '<div class="missing">Preview omitted.</div>'
    )
    marks_html = "\n".join(_render_annotation_mark(annotation, index=index) for index, annotation in enumerate(annotations, start=1))
    list_html = "\n".join(_render_annotation_list_item(annotation, index=index) for index, annotation in enumerate(annotations, start=1))
    if not list_html:
        list_html = '<div class="annotation"><strong>No annotations</strong><p>No issues were mapped to this page.</p></div>'
    return f"""<article class="panel">
  <h2>{name}</h2>
  <div class="page-grid">
    <div class="stage">{image_html}{marks_html}</div>
    <div>{list_html}</div>
  </div>
</article>"""


def _render_annotation_mark(annotation: object, *, index: int) -> str:
    if not isinstance(annotation, PresentationAnnotation):
        return ""
    region = annotation.region
    x = _css_percent(region.get("x"), fallback=0.08)
    y = _css_percent(region.get("y"), fallback=0.08)
    width = _css_percent(region.get("width"), fallback=0.30)
    height = _css_percent(region.get("height"), fallback=0.10)
    severity = "error" if annotation.severity == "error" else "warning"
    return f"""<div class="mark {severity}" style="left:{x}; top:{y}; width:{width}; height:{height};">
  <span class="mark-label">{index}</span>
</div>"""


def _render_annotation_list_item(annotation: object, *, index: int) -> str:
    if not isinstance(annotation, PresentationAnnotation):
        return ""
    recommendation = f"<p>{html.escape(annotation.recommendation)}</p>" if annotation.recommendation else ""
    return f"""<div class="annotation">
  <strong>{index}. {html.escape(annotation.severity)}: {html.escape(annotation.code)}</strong>
  <p>{html.escape(annotation.message)}</p>
  {recommendation}
</div>"""


def _css_percent(value: object, *, fallback: float) -> str:
    number = _safe_float(value)
    if number is None:
        number = fallback
    return f"{max(0.0, min(number, 1.0)) * 100:.3f}%"


def _canvas_for_format(format: str) -> dict[str, object]:
    normalized = format.strip().lower()
    canvases = {
        "ppt169": {"format": "ppt169", "viewBox": "0 0 1280 720", "width": 1280, "height": 720},
        "ppt43": {"format": "ppt43", "viewBox": "0 0 1024 768", "width": 1024, "height": 768},
        "square": {"format": "square", "viewBox": "0 0 1080 1080", "width": 1080, "height": 1080},
        "story": {"format": "story", "viewBox": "0 0 1080 1920", "width": 1080, "height": 1920},
        "xhs": {"format": "xhs", "viewBox": "0 0 1242 1660", "width": 1242, "height": 1660},
    }
    return dict(canvases.get(normalized, canvases["ppt169"]))


def _normalize_hex_color(value: str | None, *, fallback: str) -> str:
    if not value:
        return fallback
    candidate = value.strip()
    if re.fullmatch(r"#[0-9A-Fa-f]{6}", candidate):
        return candidate.upper()
    return fallback


def _layout_for_slide(*, index: int, rhythm: str, bullet_count: int) -> str:
    if index == 1:
        return "cover"
    if rhythm == "breathing":
        return "single-message"
    if bullet_count >= 5:
        return "dense-list"
    return "title-and-bullets"


def _preflight_pptx(path: Path, *, contract: PresentationContract | None = None) -> PresentationPreflightResult:
    issues: list[PresentationPreflightIssue] = []
    if not path.exists():
        issues.append(PresentationPreflightIssue("error", "missing_file", "presentation file does not exist", str(path)))
        return PresentationPreflightResult("failed", "pptx", (str(path),), tuple(issues))
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise ValueError("pptx preflight requires python-pptx") from exc
    try:
        deck = Presentation(str(path))
    except Exception as exc:  # noqa: BLE001
        issues.append(PresentationPreflightIssue("error", "invalid_pptx", f"unable to open pptx: {exc}", str(path)))
        return PresentationPreflightResult("failed", "pptx", (str(path),), tuple(issues))
    slide_count = len(deck.slides)
    if slide_count == 0:
        issues.append(PresentationPreflightIssue("error", "empty_deck", "pptx contains no slides", str(path)))
    if contract is not None and slide_count < len(contract.slides):
        issues.append(
            PresentationPreflightIssue(
                "warning",
                "slide_count_below_contract",
                f"pptx has {slide_count} slide(s), contract describes {len(contract.slides)} content slide(s)",
                str(path),
            )
        )
    slide_titles = []
    shape_count = 0
    picture_count = 0
    for index, slide in enumerate(deck.slides, start=1):
        title_text = ""
        if slide.shapes.title is not None:
            title_text = (slide.shapes.title.text or "").strip()
        if not title_text:
            title_text = _first_slide_text(slide)
        if not title_text:
            issues.append(PresentationPreflightIssue("warning", "missing_slide_title", f"slide {index} has no title", str(path)))
        slide_titles.append(title_text or f"Slide {index}")
        for shape in slide.shapes:
            shape_type = str(getattr(shape, "shape_type", ""))
            if "PICTURE" in shape_type:
                picture_count += 1
            else:
                shape_count += 1
    status = "failed" if any(issue.severity == "error" for issue in issues) else "passed"
    return PresentationPreflightResult(
        status,
        "pptx",
        (str(path),),
        tuple(issues),
        summary={
            "slide_count": slide_count,
            "titles": slide_titles[:20],
            "shape_count": shape_count,
            "picture_count": picture_count,
        },
    )


def _first_slide_text(slide) -> str:
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = (getattr(shape, "text", "") or "").strip()
        if not text:
            continue
        if re.fullmatch(r"\d{1,3}", text):
            continue
        if "/" in text and len(text) <= 32:
            continue
        return text.splitlines()[0].strip()
    return ""


def _preflight_svg(path: Path, *, contract: PresentationContract | None = None) -> PresentationPreflightResult:
    issues: list[PresentationPreflightIssue] = []
    if not path.exists():
        issues.append(PresentationPreflightIssue("error", "missing_file", "svg file does not exist", str(path)))
        return PresentationPreflightResult("failed", "svg", (str(path),), tuple(issues))
    content = path.read_text(encoding="utf-8")
    try:
        root = ET.fromstring(content)
    except ET.ParseError as exc:
        issues.append(PresentationPreflightIssue("error", "invalid_xml", f"svg is not well-formed XML: {exc}", str(path)))
        return PresentationPreflightResult("failed", "svg", (str(path),), tuple(issues))
    forbidden_patterns = {
        "foreignObject": "<foreignObject",
        "style": "<style",
        "script": "<script",
        "iframe": "<iframe",
        "textPath": "<textPath",
        "animate": "<animate",
        "class": " class=",
        "rgba": "rgba(",
        "group_opacity": "<g opacity=",
    }
    for code, needle in forbidden_patterns.items():
        if needle in content:
            issues.append(PresentationPreflightIssue("error", f"svg_forbidden_{code}", f"SVG uses unsupported feature: {needle}", str(path)))
    viewbox = root.attrib.get("viewBox") or root.attrib.get("viewbox")
    expected_viewbox = str((contract.canvas if contract is not None else _canvas_for_format("ppt169")).get("viewBox"))
    if not viewbox:
        issues.append(PresentationPreflightIssue("error", "svg_missing_viewbox", "SVG is missing viewBox", str(path)))
    elif contract is not None and viewbox.strip() != expected_viewbox:
        issues.append(
            PresentationPreflightIssue(
                "warning",
                "svg_viewbox_contract_mismatch",
                f"SVG viewBox '{viewbox}' does not match contract '{expected_viewbox}'",
                str(path),
            )
        )
    hex_colors = sorted(set(re.findall(r"#[0-9A-Fa-f]{3,8}", content)))
    status = "failed" if any(issue.severity == "error" for issue in issues) else "passed"
    return PresentationPreflightResult(
        status,
        "svg",
        (str(path),),
        tuple(issues),
        summary={"viewBox": viewbox, "color_count": len(hex_colors), "colors": hex_colors[:20]},
    )
