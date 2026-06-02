from __future__ import annotations

from pathlib import Path
import re
import shutil
import subprocess
from typing import Any

from .contracts import (
    DocumentCompanion,
    DocumentExtractionInfo,
    DocumentIngestionResult,
    ExtractedDocumentResult,
    MARKDOWN_COMPANION_KIND,
)


CONVERTIBLE_EXTENSIONS = {
    ".doc",
    ".docx",
    ".pdf",
    ".ppt",
    ".pptx",
    ".xls",
    ".xlsx",
}

_MIN_CHARS_PER_PAGE = 50
_DEFAULT_MAX_OUTLINE_ENTRIES = 50
_DEFAULT_PREVIEW_LINE_COUNT = 5


def is_convertible_document(path: Path) -> bool:
    return path.suffix.lower() in CONVERTIBLE_EXTENSIONS


def ingest_document(
    file_path: Path,
    *,
    convert_documents: bool = True,
    pdf_converter: str = "auto",
    ocr_enabled: bool = True,
    ocr_strategy: str = "local",
    ocr_languages: str = "eng+chi_sim",
    max_ocr_pages: int = 20,
    max_outline_entries: int = _DEFAULT_MAX_OUTLINE_ENTRIES,
    preview_line_count: int = _DEFAULT_PREVIEW_LINE_COUNT,
) -> DocumentIngestionResult:
    extension = file_path.suffix.lower()
    if not convert_documents or extension not in CONVERTIBLE_EXTENSIONS:
        return DocumentIngestionResult(
            extension=extension,
            extraction=DocumentExtractionInfo(status="skipped"),
        )

    markdown_text = ""
    provider: str | None = None
    ocr_provider: str | None = None
    diagnostics: list[str] = []
    page_count: int | None = None
    text_layer_present: bool | None = None
    conversion_error: str | None = None

    try:
        if extension == ".pdf":
            markdown_text, provider, ocr_provider, page_count, text_layer_present, diagnostics = _convert_pdf(
                file_path=file_path,
                pdf_converter=pdf_converter,
                ocr_enabled=ocr_enabled,
                ocr_strategy=ocr_strategy,
                ocr_languages=ocr_languages,
                max_ocr_pages=max_ocr_pages,
            )
        else:
            markdown_text, provider, diagnostics = _convert_office_document(file_path)
    except Exception as exc:  # noqa: BLE001
        conversion_error = str(exc)
        diagnostics.append(str(exc))

    if not markdown_text.strip():
        if not conversion_error:
            conversion_error = "No analysis companion was generated."
        extraction = DocumentExtractionInfo(
            status="failed",
            provider=provider,
            ocr_provider=ocr_provider,
            page_count=page_count,
            text_layer_present=text_layer_present,
            diagnostics=tuple(diagnostics or [conversion_error]),
        )
        return DocumentIngestionResult(
            extension=extension,
            extraction=extraction,
            conversion_error=conversion_error,
        )

    markdown_path = file_path.with_suffix(".md")
    markdown_path.write_text(markdown_text, encoding="utf-8")
    outline = extract_outline(markdown_path, max_entries=max_outline_entries)
    preview = [] if outline else extract_preview(markdown_path, max_lines=preview_line_count)
    extraction = DocumentExtractionInfo(
        status="completed",
        provider=provider,
        ocr_provider=ocr_provider,
        page_count=page_count,
        text_layer_present=text_layer_present,
        diagnostics=tuple(diagnostics),
    )
    companions = (
        DocumentCompanion(
            kind=MARKDOWN_COMPANION_KIND,
            label=markdown_path.name,
            path=markdown_path,
            provider=provider,
            internal=False,
        ),
    )
    return DocumentIngestionResult(
        extension=extension,
        markdown_path=markdown_path,
        companions=companions,
        extraction=extraction,
        outline=outline,
        outline_preview=preview,
        conversion_error=conversion_error,
    )


def extract_document(
    file_path: Path,
    *,
    prefer_companion: bool = True,
    convert_documents: bool = True,
    pdf_converter: str = "auto",
    ocr_enabled: bool = True,
    ocr_strategy: str = "local",
    ocr_languages: str = "eng+chi_sim",
    max_ocr_pages: int = 20,
) -> ExtractedDocumentResult:
    if prefer_companion:
        companion_path = file_path.with_suffix(".md")
        if companion_path.exists():
            return ExtractedDocumentResult(
                source_path=file_path,
                content_path=companion_path,
                content=companion_path.read_text(encoding="utf-8"),
                extraction=DocumentExtractionInfo(status="completed", provider=MARKDOWN_COMPANION_KIND),
                companions=(
                    DocumentCompanion(
                        kind=MARKDOWN_COMPANION_KIND,
                        label=companion_path.name,
                        path=companion_path,
                        provider=MARKDOWN_COMPANION_KIND,
                    ),
                ),
            )

    if file_path.suffix.lower() == ".pdf":
        text, provider, ocr_provider, page_count, text_layer_present, diagnostics = _convert_pdf(
            file_path=file_path,
            pdf_converter=pdf_converter,
            ocr_enabled=ocr_enabled,
            ocr_strategy=ocr_strategy,
            ocr_languages=ocr_languages,
            max_ocr_pages=max_ocr_pages,
        )
        if not text.strip():
            raise ValueError("No text could be extracted from the PDF.")
        extraction = DocumentExtractionInfo(
            status="completed",
            provider=provider,
            ocr_provider=ocr_provider,
            page_count=page_count,
            text_layer_present=text_layer_present,
            diagnostics=tuple(diagnostics),
        )
        return ExtractedDocumentResult(
            source_path=file_path,
            content_path=file_path,
            content=text,
            extraction=extraction,
        )

    if is_convertible_document(file_path):
        text, provider, diagnostics = _convert_office_document(file_path)
        if not text.strip():
            raise ValueError("No text could be extracted from the document.")
        extraction = DocumentExtractionInfo(
            status="completed",
            provider=provider,
            diagnostics=tuple(diagnostics),
        )
        return ExtractedDocumentResult(
            source_path=file_path,
            content_path=file_path,
            content=text,
            extraction=extraction,
        )

    return ExtractedDocumentResult(
        source_path=file_path,
        content_path=file_path,
        content=file_path.read_text(encoding="utf-8"),
        extraction=DocumentExtractionInfo(status="completed", provider="text"),
    )


def extract_outline(md_path: Path, *, max_entries: int = _DEFAULT_MAX_OUTLINE_ENTRIES) -> list[dict[str, Any]]:
    outline: list[dict[str, Any]] = []
    heading_re = re.compile(r"^(#+)\s+(.*)$")
    split_bold_heading_re = re.compile(r"^\*\*[\dA-Z][\d\.]*\*\*\s+\*\*[^*]+\*\*(?:\s+\*\*[^*]+\*\*){0,2}\s*$")
    bold_heading_re = re.compile(r"^\*\*((ITEM|PART|SECTION|SCHEDULE|EXHIBIT|APPENDIX|ANNEX|CHAPTER)\b[A-Z0-9 .,\-]*)\*\*\s*$")
    try:
        lines = md_path.read_text(encoding="utf-8").splitlines()
    except Exception:
        return []

    for index, raw_line in enumerate(lines, start=1):
        stripped = raw_line.strip()
        if not stripped:
            continue
        title: str | None = None
        if heading_match := heading_re.match(stripped):
            title = _clean_bold_title(heading_match.group(2))
        elif bold_match := bold_heading_re.match(stripped):
            title = bold_match.group(1).strip()
        elif split_bold_heading_re.match(stripped):
            title = " ".join(re.findall(r"\*\*([^*]+)\*\*", stripped))
        if title:
            outline.append({"title": title, "line": index})
        if len(outline) >= max_entries:
            outline.append({"truncated": True})
            break
    return outline


def extract_preview(md_path: Path, *, max_lines: int = _DEFAULT_PREVIEW_LINE_COUNT) -> list[str]:
    preview: list[str] = []
    try:
        with md_path.open(encoding="utf-8") as handle:
            for line in handle:
                stripped = line.strip()
                if stripped:
                    preview.append(stripped)
                if len(preview) >= max_lines:
                    break
    except Exception:
        return []
    return preview


def extract_pdf_text(
    file_path: Path,
    *,
    pdf_converter: str = "auto",
    ocr_enabled: bool = True,
    ocr_strategy: str = "local",
    ocr_languages: str = "eng+chi_sim",
    max_ocr_pages: int = 20,
) -> tuple[str | None, str | None, bool, str | None]:
    try:
        text, provider, ocr_provider, _page_count, _text_layer_present, _diagnostics = _convert_pdf(
            file_path=file_path,
            pdf_converter=pdf_converter,
            ocr_enabled=ocr_enabled,
            ocr_strategy=ocr_strategy,
            ocr_languages=ocr_languages,
            max_ocr_pages=max_ocr_pages,
        )
        return text, provider, bool(ocr_provider), None
    except Exception as exc:  # noqa: BLE001
        return None, None, False, str(exc)


def _convert_office_document(file_path: Path) -> tuple[str, str, list[str]]:
    extension = file_path.suffix.lower()
    diagnostics: list[str] = []
    if extension in {".doc", ".docx"}:
        text = _convert_docx_with_python_docx(file_path)
        if text.strip():
            return text, "python-docx", diagnostics
        diagnostics.append("python-docx returned no text.")
    elif extension in {".ppt", ".pptx"}:
        text = _convert_pptx_with_python_pptx(file_path)
        if text.strip():
            return text, "python-pptx", diagnostics
        diagnostics.append("python-pptx returned no text.")
    elif extension in {".xls", ".xlsx"}:
        text = _convert_xlsx_with_openpyxl(file_path)
        if text.strip():
            return text, "openpyxl", diagnostics
        diagnostics.append("openpyxl returned no text.")

    markitdown_text = _convert_with_markitdown(file_path)
    if markitdown_text.strip():
        return markitdown_text, "markitdown", diagnostics
    diagnostics.append("markitdown returned no text.")
    return "", "markitdown", diagnostics


def _convert_docx_with_python_docx(file_path: Path) -> str:
    try:
        from docx import Document
    except Exception:
        return ""

    try:
        document = Document(str(file_path))
    except Exception:
        return ""

    lines: list[str] = []
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        style_name = getattr(paragraph.style, "name", "") or ""
        heading_level = _heading_level_from_style(style_name)
        if heading_level > 0:
            lines.append(f"{'#' * heading_level} {text}")
        else:
            lines.append(text)

    for table in document.tables:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            lines.append("")
            lines.extend(rows)

    return "\n\n".join(lines)


def _convert_pptx_with_python_pptx(file_path: Path) -> str:
    try:
        from pptx import Presentation
    except Exception:
        return ""

    try:
        presentation = Presentation(str(file_path))
    except Exception:
        return ""

    lines: list[str] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_lines: list[str] = [f"# Slide {slide_index}"]
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                text = "".join(run.text for run in paragraph.runs).strip()
                if not text:
                    continue
                if paragraph.level == 0 and not slide_lines[-1].startswith("## "):
                    slide_lines.append(f"## {text}")
                else:
                    indent = "  " * max(paragraph.level, 0)
                    slide_lines.append(f"{indent}- {text}")
        if len(slide_lines) > 1:
            lines.append("\n".join(slide_lines))
    return "\n\n".join(lines)


def _convert_xlsx_with_openpyxl(file_path: Path) -> str:
    try:
        from openpyxl import load_workbook
    except Exception:
        return ""

    try:
        workbook = load_workbook(filename=str(file_path), data_only=True)
    except Exception:
        return ""

    sections: list[str] = []
    for sheet in workbook.worksheets:
        rows: list[str] = [f"# {sheet.title}"]
        for row in sheet.iter_rows(values_only=True):
            values = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
            if values:
                rows.append(" | ".join(values))
        if len(rows) > 1:
            sections.append("\n".join(rows))
    return "\n\n".join(sections)


def _convert_pdf(
    *,
    file_path: Path,
    pdf_converter: str,
    ocr_enabled: bool,
    ocr_strategy: str,
    ocr_languages: str,
    max_ocr_pages: int,
) -> tuple[str, str | None, str | None, int | None, bool | None, list[str]]:
    diagnostics: list[str] = []
    normalized_pdf_converter = pdf_converter.strip().lower()
    if normalized_pdf_converter not in {"auto", "pymupdf4llm", "markitdown"}:
        normalized_pdf_converter = "auto"

    page_count = _pdf_page_count(file_path)
    text_layer_present = _pdf_has_text_layer(file_path)

    if normalized_pdf_converter != "markitdown":
        pymupdf_text = _convert_pdf_with_pymupdf4llm(file_path)
        if pymupdf_text:
            if normalized_pdf_converter == "pymupdf4llm" or not _pdf_output_too_sparse(pymupdf_text, page_count):
                return pymupdf_text, "pymupdf4llm", None, page_count, text_layer_present, diagnostics
            diagnostics.append("pymupdf4llm output was too sparse; falling back.")
        else:
            diagnostics.append("pymupdf4llm unavailable or produced no text.")

    try:
        markitdown_text = _convert_with_markitdown(file_path)
        if markitdown_text.strip():
            return markitdown_text, "markitdown", None, page_count, text_layer_present, diagnostics
        diagnostics.append("markitdown produced no text.")
    except Exception as exc:
        diagnostics.append(f"markitdown failed: {exc}")
        if normalized_pdf_converter == "markitdown" and not ocr_enabled:
            raise

    marker_text = _convert_pdf_with_marker_pdf(file_path)
    if marker_text:
        diagnostics.append("marker-pdf fallback succeeded.")
        return marker_text, "marker-pdf", "marker-pdf", page_count, text_layer_present, diagnostics
    diagnostics.append("marker-pdf unavailable or produced no text.")

    pypdf_text = _convert_pdf_with_pypdf2(file_path)
    if pypdf_text.strip():
        return pypdf_text, "pypdf2", None, page_count, text_layer_present, diagnostics
    diagnostics.append("pypdf2 produced no text.")

    if ocr_enabled and ocr_strategy == "local":
        ocr_text = _convert_pdf_with_local_ocr(
            file_path=file_path,
            languages=ocr_languages,
            max_pages=max_ocr_pages,
        )
        if ocr_text.strip():
            return ocr_text, "tesseract", "tesseract", page_count, text_layer_present, diagnostics
        diagnostics.append("tesseract OCR produced no text.")

    raise ValueError("No text could be extracted from the PDF.")


def _convert_with_markitdown(file_path: Path) -> str:
    from markitdown import MarkItDown

    return MarkItDown().convert(str(file_path)).text_content


def _convert_pdf_with_pymupdf4llm(file_path: Path) -> str | None:
    try:
        import pymupdf4llm
    except Exception:
        return None

    try:
        return pymupdf4llm.to_markdown(str(file_path))
    except Exception:
        return None


def _convert_pdf_with_marker_pdf(file_path: Path) -> str | None:
    marker_single = shutil.which("marker_single")
    if marker_single is None:
        return None
    output_root = file_path.parent / ".marker-output"
    output_root.mkdir(parents=True, exist_ok=True)
    command = [marker_single, str(file_path), "--output_dir", str(output_root)]
    completed = subprocess.run(command, capture_output=True, text=True, check=False)
    if completed.returncode != 0:
        return None
    markdown_path = output_root / file_path.stem / f"{file_path.stem}.md"
    if markdown_path.exists():
        try:
            return markdown_path.read_text(encoding="utf-8")
        except Exception:
            return None
    return None


def _convert_pdf_with_pypdf2(file_path: Path) -> str:
    try:
        from PyPDF2 import PdfReader
    except Exception:
        return ""

    try:
        reader = PdfReader(str(file_path))
    except Exception:
        return ""

    chunks: list[str] = []
    for page_index, page in enumerate(reader.pages, start=1):
        try:
            text = (page.extract_text() or "").strip()
        except Exception:
            text = ""
        if text:
            chunks.append(f"[Page {page_index}]\n{text}")
    return "\n\n".join(chunks)


def _convert_pdf_with_local_ocr(
    *,
    file_path: Path,
    languages: str,
    max_pages: int,
) -> str:
    try:
        import fitz
        import pytesseract
        from PIL import Image
    except Exception:
        return ""

    document = fitz.open(str(file_path))
    chunks: list[str] = []
    try:
        for page_index, page in enumerate(document, start=1):
            if page_index > max_pages:
                chunks.append(f"[OCR truncated after {max_pages} pages]")
                break
            pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
            image = Image.frombytes("RGB", [pixmap.width, pixmap.height], pixmap.samples)
            try:
                text = pytesseract.image_to_string(image, lang=languages).strip()
            except Exception:
                return ""
            if text:
                chunks.append(f"[Page {page_index}]\n{text}")
    finally:
        document.close()
    return "\n\n".join(chunks)


def _pdf_page_count(file_path: Path) -> int | None:
    try:
        import fitz
    except Exception:
        return None

    document = fitz.open(str(file_path))
    try:
        return len(document)
    finally:
        document.close()


def _pdf_has_text_layer(file_path: Path) -> bool | None:
    try:
        import fitz
    except Exception:
        return None

    document = fitz.open(str(file_path))
    try:
        for page in document:
            if page.get_text("text").strip():
                return True
        return False
    finally:
        document.close()


def _pdf_output_too_sparse(text: str, page_count: int | None) -> bool:
    characters = len(text.strip())
    if page_count and page_count > 0:
        return (characters / page_count) < _MIN_CHARS_PER_PAGE
    return characters < 200


def _clean_bold_title(value: str) -> str:
    merged = re.sub(r"\*\*\s*\*\*", " ", value).strip()
    if match := re.fullmatch(r"\*\*(.+?)\*\*", merged, re.DOTALL):
        return match.group(1).strip()
    return merged


def _heading_level_from_style(style_name: str) -> int:
    lowered = style_name.lower()
    if not lowered.startswith("heading"):
        return 0
    match = re.search(r"(\d+)", lowered)
    if not match:
        return 1
    return max(1, min(int(match.group(1)), 6))
